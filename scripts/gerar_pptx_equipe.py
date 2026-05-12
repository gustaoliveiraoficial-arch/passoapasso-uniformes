# -*- coding: utf-8 -*-
import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Paleta ──────────────────────────────────────────────────────────────────
PRETO      = RGBColor(0x06, 0x06, 0x06)
ESCURO     = RGBColor(0x0E, 0x0E, 0x0E)
CARD       = RGBColor(0x18, 0x18, 0x18)
BORDA      = RGBColor(0x2A, 0x2A, 0x2A)
OURO       = RGBColor(0xF5, 0x9E, 0x0B)
OURO_CLARO = RGBColor(0xFC, 0xD3, 0x4D)
LARANJA    = RGBColor(0xF9, 0x73, 0x16)
VERMELHO   = RGBColor(0xDC, 0x26, 0x26)
VERDE      = RGBColor(0x10, 0xB9, 0x81)
AZUL       = RGBColor(0x60, 0xA5, 0xFA)
ROXO       = RGBColor(0xA7, 0x8B, 0xFA)
ROSA       = RGBColor(0xF4, 0x72, 0xB6)
BRANCO     = RGBColor(0xFF, 0xFF, 0xFF)
CINZA      = RGBColor(0x6B, 0x72, 0x80)
CINZA_CLR  = RGBColor(0xD1, 0xD5, 0xDB)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ─── helpers ────────────────────────────────────────────────────────────────

def bg(slide, color=PRETO):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def rect(slide, l, t, w, h, fill=None, line=None, lw=Pt(0.6)):
    s = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.width = lw
    else:
        s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=14, bold=False, color=BRANCO,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.italic = italic
    return tb

def bar(slide, l, t, w=Inches(0.5), h=Pt(3), color=OURO):
    rect(slide, l, t, w, h, fill=color)

def num(slide, n):
    txt(slide, str(n).zfill(2), Inches(12.8), Inches(7.1), Inches(0.5),
        Inches(0.35), size=9, color=CINZA, align=PP_ALIGN.RIGHT)

def card(slide, l, t, w, h, fill=CARD, line=BORDA, lw=Pt(0.5), accent=None):
    rect(slide, l, t, w, h, fill=fill, line=line, lw=lw)
    if accent:
        rect(slide, l, t, Pt(3), h, fill=accent)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — CAPA EMOCIONAL
# ─────────────────────────────────────────────────────────────────────────────
def s1():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, PRETO)
    # Faixa dourada superior
    rect(slide, Inches(0), Inches(0), W, Pt(4), fill=OURO)
    # Faixa dourada inferior
    rect(slide, Inches(0), H - Pt(4), W, Pt(4), fill=OURO)
    # Bloco escuro lateral
    rect(slide, Inches(0), Inches(0), Inches(0.12), H, fill=RGBColor(0x2A,0x1A,0x00))

    txt(slide, "PASSO A PASSO UNIFORMES",
        Inches(0.6), Inches(0.35), Inches(12), Inches(0.4),
        size=9, color=CINZA, align=PP_ALIGN.CENTER)

    txt(slide, "Uma nova era",
        Inches(0.6), Inches(1.2), Inches(12), Inches(0.9),
        size=48, bold=False, color=CINZA_CLR, align=PP_ALIGN.CENTER)

    txt(slide, "comeca hoje.",
        Inches(0.6), Inches(2.1), Inches(12), Inches(1.1),
        size=72, bold=True, color=OURO, align=PP_ALIGN.CENTER)

    bar(slide, Inches(6.1), Inches(3.35), Inches(1.2), color=OURO)

    txt(slide, "Esta reuniao vai mudar a forma como voce enxerga\no seu trabalho, o seu potencial e o seu futuro.",
        Inches(2.0), Inches(3.55), Inches(9.3), Inches(0.9),
        size=15, color=CINZA_CLR, align=PP_ALIGN.CENTER, italic=True)

    txt(slide, "Abril 2026  |  Equipe de Vendas  |  Passo a Passo Uniformes",
        Inches(2.0), Inches(6.9), Inches(9.3), Inches(0.4),
        size=9, color=CINZA, align=PP_ALIGN.CENTER)
    num(slide, 1)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — POR QUE VOCE ESTA AQUI
# ─────────────────────────────────────────────────────────────────────────────
def s2():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, PRETO)
    rect(slide, Inches(0), Inches(0), Inches(0.06), H, fill=OURO)

    txt(slide, "Voce nao esta aqui por acaso.",
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.75),
        size=34, bold=True, color=OURO)

    bar(slide, Inches(0.5), Inches(1.35))
    txt(slide, "Cada pessoa nesta sala foi escolhida porque tem potencial para algo maior.",
        Inches(0.5), Inches(1.5), Inches(9.5), Inches(0.55),
        size=14, color=CINZA_CLR)

    perguntas = [
        ("Voce acorda pensando em como pode ser melhor hoje?"),
        ("Voce sente que o seu trabalho pode mudar a vida da sua familia?"),
        ("Voce acredita que tem mais a dar do que esta dando?"),
        ("Voce quer ser reconhecido, crescer, e se orgulhar do que faz?"),
    ]
    for i, p in enumerate(perguntas):
        y = Inches(2.2) + i * Inches(1.0)
        card(slide, Inches(0.5), y, Inches(8.5), Inches(0.82))
        txt(slide, "?", Inches(0.65), y + Pt(6), Inches(0.4), Inches(0.65),
            size=22, bold=True, color=OURO)
        txt(slide, p, Inches(1.2), y + Pt(8), Inches(7.6), Inches(0.6),
            size=13, color=CINZA_CLR)

    # Box resposta
    rect(slide, Inches(0.5), Inches(6.4), Inches(8.5), Inches(0.72),
         fill=RGBColor(0x1A,0x12,0x00), line=RGBColor(0x78,0x50,0x1A), lw=Pt(0.75))
    rect(slide, Inches(0.5), Inches(6.4), Pt(3), Inches(0.72), fill=OURO)
    txt(slide, "Se voce respondeu SIM para qualquer dessas perguntas — esta reuniao e para voce.",
        Inches(0.65), Inches(6.48), Inches(8.2), Inches(0.55),
        size=11.5, bold=True, color=OURO)

    # Lado direito — citacao
    rect(slide, Inches(9.3), Inches(0.5), Inches(3.8), Inches(6.7),
         fill=RGBColor(0x12,0x0D,0x00), line=RGBColor(0x3A,0x28,0x00), lw=Pt(0.5))
    txt(slide, '"',
        Inches(9.5), Inches(0.6), Inches(3.5), Inches(1.2),
        size=80, bold=True, color=OURO, align=PP_ALIGN.CENTER)
    txt(slide, "O sucesso nao e um destino. E uma identidade que voce escolhe todos os dias.",
        Inches(9.5), Inches(1.8), Inches(3.4), Inches(2.0),
        size=14, italic=True, color=CINZA_CLR, align=PP_ALIGN.CENTER)
    txt(slide, "— James Clear\nAutor de Atomic Habits",
        Inches(9.5), Inches(4.2), Inches(3.4), Inches(0.8),
        size=10, color=CINZA, align=PP_ALIGN.CENTER)
    num(slide, 2)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — O SEU SONHO
# ─────────────────────────────────────────────────────────────────────────────
def s3():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, RGBColor(0x06,0x04,0x01))
    rect(slide, Inches(0), Inches(0), W, Pt(4), fill=OURO)

    txt(slide, "Feche os olhos por um segundo.",
        Inches(0.6), Inches(0.35), Inches(12), Inches(0.6),
        size=13, italic=True, color=CINZA, align=PP_ALIGN.CENTER)

    txt(slide, "O que voce esta",
        Inches(0.6), Inches(1.1), Inches(12), Inches(0.75),
        size=40, bold=False, color=CINZA_CLR, align=PP_ALIGN.CENTER)
    txt(slide, "correndo atras?",
        Inches(0.6), Inches(1.85), Inches(12), Inches(0.85),
        size=52, bold=True, color=OURO, align=PP_ALIGN.CENTER)

    bar(slide, Inches(6.1), Inches(2.82), Inches(1.2), color=LARANJA)

    sonhos = [
        ("🏠", "Uma casa propria.\nA sensacao de ter algo seu."),
        ("🚗", "Um carro novo.\nNao precisar pedir carona."),
        ("👨‍👩‍👧", "Dar para sua familia\no que ela merece."),
        ("📚", "Pagar a faculdade\ndos seus filhos."),
        ("✈️", "Uma viagem que voce\nnunca imaginou fazer."),
        ("💪", "Provar para si mesmo\nque e capaz."),
    ]

    for i, (emoji, texto) in enumerate(sonhos):
        col = i % 3
        row = i // 3
        x = Inches(0.5) + col * Inches(4.28)
        y = Inches(3.1) + row * Inches(1.85)
        rect(slide, x, y, Inches(4.0), Inches(1.65),
             fill=RGBColor(0x14,0x0F,0x02), line=RGBColor(0x3A,0x28,0x00), lw=Pt(0.5))
        rect(slide, x, y, Pt(3), Inches(1.65), fill=OURO)
        txt(slide, emoji, x + Inches(0.2), y + Pt(8), Inches(0.6), Inches(1.0), size=28)
        txt(slide, texto, x + Inches(0.9), y + Pt(12), Inches(2.9), Inches(1.3),
            size=12.5, color=CINZA_CLR)

    num(slide, 3)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — NEUROCIENCIA: SEU CEREBRO E UMA MAQUINA
# ─────────────────────────────────────────────────────────────────────────────
def s4():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, PRETO)
    rect(slide, Inches(0), Inches(0), Inches(0.06), H, fill=ROXO)

    txt(slide, "NEUROCIENCIA",
        Inches(0.5), Inches(0.3), Inches(8), Inches(0.35),
        size=9, bold=True, color=ROXO)
    bar(slide, Inches(0.5), Inches(0.72), color=ROXO)
    txt(slide, "Seu cerebro esta do seu lado.",
        Inches(0.5), Inches(0.85), Inches(9), Inches(0.7),
        size=32, bold=True, color=BRANCO)
    txt(slide, "Voce so precisa entender como ele funciona.",
        Inches(0.5), Inches(1.55), Inches(9), Inches(0.4),
        size=13, color=CINZA_CLR)

    neuro = [
        (ROXO, "DOPAMINA",
         "Seu cerebro libera dopamina ANTES da recompensa — nao depois.\n"
         "O simples ato de ligar para um cliente e seguir o processo\n"
         "ja ativa o sistema de recompensa. Voce precisa comecar para sentir."),
        (AZUL, "NEUROPLASTICIDADE",
         "Cada vez que voce repete uma acao, seu cerebro cria\n"
         "novas conexoes neuronais. Vendas e um musculo.\n"
         "Quanto mais voce treina, mais facil fica. Sempre."),
        (VERDE, "NEURONIOS ESPELHO",
         "Seu cerebro copia as pessoas ao redor.\n"
         "Se voce se cerca de pessoas que vendem bem e tem alto padrao,\n"
         "seu cerebro automaticamente eleva o seu padrao tambem."),
        (OURO, "VISUALIZACAO",
         "Estudos da Harvard mostram que o cerebro nao distingue\n"
         "entre uma experiencia real e uma visualizada com detalhes.\n"
         "Visualizar o sucesso ja esta te preparando para ele."),
    ]

    for i, (cor, titulo, descr) in enumerate(neuro):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * Inches(6.4)
        y = Inches(2.1) + row * Inches(2.35)
        card(slide, x, y, Inches(6.1), Inches(2.1), accent=cor)
        txt(slide, titulo, x + Inches(0.2), y + Pt(8), Inches(5.7), Inches(0.38),
            size=10, bold=True, color=cor)
        txt(slide, descr, x + Inches(0.2), y + Pt(30), Inches(5.7), Inches(1.5),
            size=11, color=CINZA_CLR)

    # Rodape
    rect(slide, Inches(0.5), Inches(6.85), Inches(12.3), Pt(1), fill=BORDA)
    txt(slide, "Fontes: Pesquisa da Harvard Medical School · James Clear (Atomic Habits) · Andrew Huberman (Huberman Lab)",
        Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.35),
        size=8, color=CINZA, align=PP_ALIGN.CENTER)
    num(slide, 4)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — IDENTIDADE: VOCE E UM LEAO
# ─────────────────────────────────────────────────────────────────────────────
def s5():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, PRETO)
    rect(slide, Inches(0), Inches(0), W, Pt(4), fill=LARANJA)

    txt(slide, "Existe uma diferenca entre fazer e SER.",
        Inches(0.6), Inches(0.3), Inches(12), Inches(0.55),
        size=13, italic=True, color=CINZA_CLR, align=PP_ALIGN.CENTER)

    txt(slide, "Voce nao \"tenta vender\".",
        Inches(0.6), Inches(1.05), Inches(12), Inches(0.7),
        size=36, bold=False, color=CINZA_CLR, align=PP_ALIGN.CENTER)
    txt(slide, "Voce E um vendedor.",
        Inches(0.6), Inches(1.75), Inches(12), Inches(0.8),
        size=48, bold=True, color=LARANJA, align=PP_ALIGN.CENTER)

    bar(slide, Inches(6.1), Inches(2.7), Inches(1.2), color=LARANJA)

    # Tabela comparativa
    comparativo = [
        ("Mentalidade Antiga",   "Nova Identidade"),
        ('"Preciso vender mais"',  '"Eu sou um profissional de vendas"'),
        ('"Espero o cliente aparecer"', '"Eu caço oportunidades"'),
        ('"Tudo bem nao fechar"', '"Cada nao me aproxima de um sim"'),
        ('"O processo e chato"',  '"O processo e o que me torna imparavel"'),
        ('"Sorte dos outros"',    '"Eu crio minha propria sorte"'),
    ]

    for i, (esq, dir_) in enumerate(comparativo):
        y = Inches(2.95) + i * Inches(0.65)
        is_header = (i == 0)
        fb = RGBColor(0x1A,0x0C,0x0C) if is_header else RGBColor(0x16,0x0A,0x0A)
        fb2 = RGBColor(0x0A,0x18,0x0A) if is_header else RGBColor(0x07,0x14,0x07)
        rect(slide, Inches(0.5), y, Inches(5.8), Inches(0.58),
             fill=fb, line=BORDA, lw=Pt(0.3))
        rect(slide, Inches(6.5), y, Inches(6.4), Inches(0.58),
             fill=fb2, line=BORDA, lw=Pt(0.3))
        # seta
        txt(slide, "->", Inches(6.2), y + Pt(6), Inches(0.35), Inches(0.4),
            size=11, bold=True, color=CINZA, align=PP_ALIGN.CENTER)
        cor_e = VERMELHO if is_header else CINZA_CLR
        cor_d = VERDE if is_header else BRANCO
        b_e = True if is_header else False
        txt(slide, esq,  Inches(0.65), y + Pt(7), Inches(5.5), Inches(0.45),
            size=11 if not is_header else 10, bold=b_e, color=cor_e)
        txt(slide, dir_, Inches(6.65), y + Pt(7), Inches(6.1), Inches(0.45),
            size=11 if not is_header else 10, bold=is_header, color=cor_d)

    num(slide, 5)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — O PODER DO PROCESSO
# ─────────────────────────────────────────────────────────────────────────────
def s6():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, PRETO)
    rect(slide, Inches(0), Inches(0), Inches(0.06), H, fill=VERDE)

    txt(slide, "CIENCIA DOS HABITOS",
        Inches(0.5), Inches(0.3), Inches(8), Inches(0.35),
        size=9, bold=True, color=VERDE)
    bar(slide, Inches(0.5), Inches(0.72), color=VERDE)
    txt(slide, "Nao confie na motivacao.",
        Inches(0.5), Inches(0.85), Inches(8), Inches(0.65),
        size=32, bold=True, color=BRANCO)
    txt(slide, "Confie no processo. A motivacao vai e vem. O processo fica.",
        Inches(0.5), Inches(1.5), Inches(9.5), Inches(0.4),
        size=13, color=CINZA_CLR)

    # O loop do habito
    txt(slide, "O LOOP DO HABITO VENCEDOR",
        Inches(0.5), Inches(2.1), Inches(6), Inches(0.3),
        size=9, bold=True, color=VERDE)

    loop = [
        (OURO,  "1. GATILHO",    "Abre o Kommo.\nVe os leads do dia."),
        (AZUL,  "2. ROTINA",     "Segue o script.\nFaz o contato."),
        (LARANJA,"3. RECOMPENSA","Fechou? O sino toca.\nNao fechou? Aprendeu."),
        (ROXO,  "4. REPETICAO",  "Faz de novo amanha.\nSeu cerebro fica mais forte."),
    ]
    for i, (cor, titulo, descr) in enumerate(loop):
        x = Inches(0.5) + i * Inches(3.2)
        card(slide, x, Inches(2.5), Inches(3.0), Inches(2.4))
        rect(slide, x, Inches(2.5), Inches(3.0), Pt(3), fill=cor)
        txt(slide, titulo, x + Inches(0.2), Inches(2.65), Inches(2.6), Inches(0.35),
            size=10, bold=True, color=cor)
        txt(slide, descr, x + Inches(0.2), Inches(3.05), Inches(2.6), Inches(1.5),
            size=12, color=CINZA_CLR)
        if i < 3:
            txt(slide, "->", x + Inches(2.8), Inches(3.4), Inches(0.5), Inches(0.4),
                size=16, bold=True, color=BORDA, align=PP_ALIGN.CENTER)

    # Citacao central
    rect(slide, Inches(0.5), Inches(5.1), Inches(12.3), Inches(1.8),
         fill=RGBColor(0x05,0x17,0x10), line=RGBColor(0x0A,0x3D,0x22), lw=Pt(0.75))
    rect(slide, Inches(0.5), Inches(5.1), Pt(3), Inches(1.8), fill=VERDE)
    txt(slide, '"Nao e o mais forte que sobrevive, nem o mais inteligente.\nE aquele que melhor se adapta ao processo."',
        Inches(0.65), Inches(5.2), Inches(11.0), Inches(1.0),
        size=15, italic=True, color=CINZA_CLR, align=PP_ALIGN.CENTER)
    txt(slide, "— Principio da Neuroplasticidade Adaptativa",
        Inches(0.65), Inches(6.15), Inches(12.0), Inches(0.35),
        size=9, color=CINZA, align=PP_ALIGN.CENTER)
    num(slide, 6)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — QUEM LIDERA VOCES (confiar no lider)
# ─────────────────────────────────────────────────────────────────────────────
def s7():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, PRETO)
    rect(slide, Inches(0), Inches(0), W, Pt(4), fill=OURO)

    txt(slide, "Por que voce pode",
        Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
        size=36, bold=False, color=CINZA_CLR, align=PP_ALIGN.CENTER)
    txt(slide, "confiar nesta lideranca.",
        Inches(0.6), Inches(1.0), Inches(12), Inches(0.8),
        size=44, bold=True, color=OURO, align=PP_ALIGN.CENTER)

    bar(slide, Inches(6.1), Inches(1.95), Inches(1.2))

    razoes = [
        ("Estou correndo o mesmo risco que voce.",
         "Nao estou pedindo para voce fazer algo que eu mesmo nao esteja fazendo.\n"
         "Cada decisao que tomo e pensando no crescimento de todos nos."),
        ("Eu conheco o terreno.",
         "Cada cliente, cada objecao, cada canal — eu estudei, testei e validei.\n"
         "O plano que estou trazendo nao e teoria. E o resultado de muito trabalho."),
        ("Meu sucesso depende do seu sucesso.",
         "Quando voce cresce, eu cresço. Quando voce vende, a empresa cresce.\n"
         "Estamos no mesmo barco. Eu nao tenho interesse em te ver falhar."),
        ("Eu vou estar aqui todos os dias.",
         "Nao sou o tipo de lider que da o plano e some.\n"
         "Estarei na trincheira com voce. Todos os dias. Sem excecao."),
    ]

    for i, (titulo, descr) in enumerate(razoes):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * Inches(6.4)
        y = Inches(2.25) + row * Inches(2.2)
        card(slide, x, y, Inches(6.1), Inches(2.0),
             fill=RGBColor(0x14,0x10,0x02),
             line=RGBColor(0x3A,0x28,0x00), lw=Pt(0.75))
        rect(slide, x, y, Pt(3), Inches(2.0), fill=OURO)
        txt(slide, titulo, x + Inches(0.2), y + Pt(10), Inches(5.7), Inches(0.4),
            size=13, bold=True, color=OURO)
        txt(slide, descr, x + Inches(0.2), y + Pt(38), Inches(5.7), Inches(1.3),
            size=11, color=CINZA_CLR)

    num(slide, 7)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — O EFEITO COMPOSTO
# ─────────────────────────────────────────────────────────────────────────────
def s8():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, PRETO)
    rect(slide, Inches(0), Inches(0), Inches(0.06), H, fill=AZUL)

    txt(slide, "EFEITO COMPOSTO",
        Inches(0.5), Inches(0.3), Inches(8), Inches(0.35),
        size=9, bold=True, color=AZUL)
    bar(slide, Inches(0.5), Inches(0.72), color=AZUL)
    txt(slide, "Pequenas acoes todos os dias",
        Inches(0.5), Inches(0.85), Inches(12), Inches(0.65),
        size=32, bold=True, color=BRANCO)
    txt(slide, "criam resultados imposssiveis de ignorar.",
        Inches(0.5), Inches(1.5), Inches(12), Inches(0.4),
        size=18, bold=False, color=AZUL)

    # Progressao visual
    passos = [
        (CINZA,   "Dia 1",   "Voce faz a acao. Parece pequeno. E."),
        (CINZA_CLR,"Semana 1","O habito comeca a se formar. Fica mais facil."),
        (OURO_CLARO,"Semana 2","Voce ja e reconhecido como o que age."),
        (OURO,    "Mes 1",   "Os resultados aparecem. O time te ve diferente."),
        (LARANJA, "Mes 3",   "Voce virou referencia. E impossivel ignorar."),
    ]

    for i, (cor, tempo, descr) in enumerate(passos):
        y = Inches(2.15) + i * Inches(0.88)
        largura = Inches(2.0) + i * Inches(1.5)
        rect(slide, Inches(0.5), y, largura, Inches(0.72), fill=cor)
        col_txt = PRETO if cor in [OURO, OURO_CLARO, LARANJA] else BRANCO
        txt(slide, tempo, Inches(0.65), y + Pt(6), Inches(1.5), Inches(0.55),
            size=11, bold=True, color=col_txt)
        txt(slide, descr, Inches(0.65) + largura, y + Pt(10), Inches(12) - largura - Inches(0.2), Inches(0.5),
            size=11, color=CINZA_CLR)

    # Citacao
    rect(slide, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.65),
         fill=RGBColor(0x0B,0x19,0x30), line=RGBColor(0x1E,0x40,0xAF), lw=Pt(0.5))
    txt(slide, '"1% melhor a cada dia = 37x melhor ao final de um ano."  — James Clear',
        Inches(0.65), Inches(6.6), Inches(12.0), Inches(0.5),
        size=11, italic=True, color=AZUL, align=PP_ALIGN.CENTER)
    num(slide, 8)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — O AMBIENTE QUE VAMOS CRIAR
# ─────────────────────────────────────────────────────────────────────────────
def s9():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, PRETO)

    txt(slide, "Voce nao vai fazer isso sozinho.",
        Inches(0.6), Inches(0.4), Inches(12), Inches(0.65),
        size=34, bold=True, color=OURO, align=PP_ALIGN.CENTER)
    bar(slide, Inches(6.1), Inches(1.12), Inches(1.2))
    txt(slide, "Estamos construindo um ambiente onde o sucesso e inevitavel.",
        Inches(0.6), Inches(1.3), Inches(12), Inches(0.4),
        size=13, color=CINZA_CLR, align=PP_ALIGN.CENTER)

    ambiente = [
        ("Treinamento toda semana",
         "30 minutos. Todo lunes. Tecnica, produto, roleplay e resultado.\nQuem treina, vende. Quem vende, cresce."),
        ("O sino e o placar",
         "Toda venda celebrada. Toda meta registrada.\nO sucesso de um inspira o outro — isso e ciencia."),
        ("Competicao saudavel",
         "Nao estamos competindo contra o mercado.\nEstamos elevando o padrao juntos, todos os dias."),
        ("Reconhecimento real",
         "Quem bate meta e reconhecido em publico.\nNao e elogio vazio — e respeito verdadeiro pelo esforco."),
        ("Feedback honesto",
         "Nao vou te deixar errar sem te avisar.\nNao vou te abandonar quando for dificil."),
        ("Cultura de crescimento",
         "Errar faz parte. O que nao pode e repetir o mesmo erro.\nCada queda e uma aula. Cada aula te torna mais forte."),
    ]

    for i, (titulo, descr) in enumerate(ambiente):
        col = i % 3
        row = i // 3
        x = Inches(0.5) + col * Inches(4.28)
        y = Inches(1.9) + row * Inches(2.3)
        card(slide, x, y, Inches(4.0), Inches(2.1))
        rect(slide, x, y, Inches(4.0), Pt(3), fill=OURO)
        txt(slide, titulo, x + Inches(0.15), y + Pt(12), Inches(3.7), Inches(0.4),
            size=12, bold=True, color=OURO)
        txt(slide, descr, x + Inches(0.15), y + Pt(38), Inches(3.7), Inches(1.5),
            size=10.5, color=CINZA_CLR)

    num(slide, 9)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — O COMPROMISSO
# ─────────────────────────────────────────────────────────────────────────────
def s10():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, PRETO)
    rect(slide, Inches(0), Inches(0), W, Pt(4), fill=LARANJA)

    txt(slide, "Um compromisso simples.",
        Inches(0.6), Inches(0.3), Inches(12), Inches(0.65),
        size=32, bold=False, color=CINZA_CLR, align=PP_ALIGN.CENTER)
    txt(slide, "Nao com a empresa. Com voce mesmo.",
        Inches(0.6), Inches(0.95), Inches(12), Inches(0.65),
        size=36, bold=True, color=LARANJA, align=PP_ALIGN.CENTER)

    bar(slide, Inches(6.1), Inches(1.72), Inches(1.2), color=LARANJA)

    promessas = [
        "Eu vou seguir o processo — mesmo quando nao estiver com vontade.",
        "Eu vou tratar cada cliente como se fosse o mais importante.",
        "Eu vou aprender com cada nao e melhorar a cada semana.",
        "Eu vou celebrar as vitorias dos meus colegas como se fossem minhas.",
        "Eu vou aparecer todos os dias — de corpo e mente presente.",
        "Eu vou acreditar no processo mesmo antes de ver os resultados.",
    ]

    for i, p in enumerate(promessas):
        y = Inches(2.05) + i * Inches(0.75)
        rect(slide, Inches(1.5), y, Inches(10.4), Inches(0.62),
             fill=CARD, line=BORDA, lw=Pt(0.4))
        txt(slide, "[ ]", Inches(1.65), y + Pt(6), Inches(0.45), Inches(0.5),
            size=12, color=LARANJA)
        txt(slide, p, Inches(2.2), y + Pt(7), Inches(9.5), Inches(0.48),
            size=12, color=CINZA_CLR)

    rect(slide, Inches(1.5), Inches(6.55), Inches(10.4), Inches(0.62),
         fill=RGBColor(0x1A,0x09,0x02), line=LARANJA, lw=Pt(0.75))
    txt(slide, "Este nao e um papel que voce assina. E uma escolha que voce faz agora.",
        Inches(1.65), Inches(6.62), Inches(10.2), Inches(0.48),
        size=11, bold=True, color=LARANJA, align=PP_ALIGN.CENTER)
    num(slide, 10)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — FECHAMENTO EMOCIONAL
# ─────────────────────────────────────────────────────────────────────────────
def s11():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, PRETO)
    rect(slide, Inches(0), Inches(0), W, Pt(4), fill=OURO)
    rect(slide, Inches(0), H - Pt(4), W, Pt(4), fill=OURO)

    txt(slide, "Daqui a 3 meses,",
        Inches(0.6), Inches(0.6), Inches(12), Inches(0.75),
        size=36, bold=False, color=CINZA_CLR, align=PP_ALIGN.CENTER)
    txt(slide, "voce vai olhar para tras",
        Inches(0.6), Inches(1.35), Inches(12), Inches(0.75),
        size=40, bold=False, color=CINZA_CLR, align=PP_ALIGN.CENTER)
    txt(slide, "e nao vai acreditar no quanto cresceu.",
        Inches(0.6), Inches(2.1), Inches(12), Inches(0.9),
        size=46, bold=True, color=OURO, align=PP_ALIGN.CENTER)

    bar(slide, Inches(6.1), Inches(3.15), Inches(1.2))

    txt(slide, "O unico requisito e que voce apareca amanha.",
        Inches(1.0), Inches(3.4), Inches(11.3), Inches(0.6),
        size=18, italic=True, color=CINZA_CLR, align=PP_ALIGN.CENTER)

    rect(slide, Inches(2.5), Inches(4.2), Inches(8.3), Inches(1.5),
         fill=RGBColor(0x12,0x0D,0x00), line=RGBColor(0x78,0x50,0x1A), lw=Pt(1))
    txt(slide, '"A jornada de mil milhas comeca com um unico passo."',
        Inches(2.7), Inches(4.35), Inches(7.9), Inches(0.7),
        size=16, italic=True, color=OURO, align=PP_ALIGN.CENTER)
    txt(slide, "— Lao Tse",
        Inches(2.7), Inches(5.0), Inches(7.9), Inches(0.35),
        size=10, color=CINZA, align=PP_ALIGN.CENTER)

    txt(slide, "Vamos juntos.  Somos leoes.",
        Inches(0.6), Inches(5.8), Inches(12), Inches(0.7),
        size=28, bold=True, color=LARANJA, align=PP_ALIGN.CENTER)

    txt(slide, "Passo a Passo Uniformes  |  Equipe de Vendas  |  Abril 2026",
        Inches(0.6), Inches(7.1), Inches(12), Inches(0.3),
        size=8, color=CINZA, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# Gerar
# ─────────────────────────────────────────────────────────────────────────────
s1(); s2(); s3(); s4(); s5(); s6(); s7(); s8(); s9(); s10(); s11()

output = r"C:\Users\gusta\OneDrive\Área de Trabalho\PassoaPasso\docs\apresentacao-equipe-motivacional-2026.pptx"
prs.save(output)
print(f"Salvo em: {output}")
