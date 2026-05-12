from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from pptx.enum.dml import MSO_THEME_COLOR
import os

# ── Cores da identidade ──────────────────────────────────────────────────────
PRETO      = RGBColor(0x0a, 0x0a, 0x0a)
ESCURO     = RGBColor(0x11, 0x18, 0x27)
CARD       = RGBColor(0x1f, 0x29, 0x37)
BORDA      = RGBColor(0x37, 0x41, 0x51)
OURO       = RGBColor(0xF5, 0x9E, 0x0B)
OURO_CLARO = RGBColor(0xFC, 0xD3, 0x4D)
LARANJA    = RGBColor(0xF9, 0x73, 0x16)
VERDE      = RGBColor(0x10, 0xB9, 0x81)
AZUL       = RGBColor(0x3B, 0x82, 0xF6)
VERMELHO   = RGBColor(0xEF, 0x44, 0x44)
BRANCO     = RGBColor(0xFF, 0xFF, 0xFF)
CINZA      = RGBColor(0x6B, 0x72, 0x80)
CINZA_CLR  = RGBColor(0xD1, 0xD5, 0xDB)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # layout em branco

# ─────────────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────────────

def bg(slide, color=PRETO):
    """Pinta o fundo do slide."""
    bg_ = slide.background
    fill = bg_.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height,
             fill_color=None, line_color=None, line_width=Pt(0.75),
             fill_alpha=None, rx=0):
    from pptx.oxml.ns import qn
    from lxml import etree
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    # fill
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    # line
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=BRANCO, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def add_accent_bar(slide, left, top, width=Inches(0.45), height=Pt(3), color=OURO):
    """Barra decorativa horizontal (divisor)."""
    add_rect(slide, left, top, width, Pt(3), fill_color=color)

def add_badge(slide, text, left, top,
              bg_color=RGBColor(0x27,0x1E,0x07),
              text_color=OURO,
              border_color=RGBColor(0x78,0x50,0x1A)):
    w = Inches(1.8); h = Inches(0.30)
    r = add_rect(slide, left, top, w, h, fill_color=bg_color, line_color=border_color, line_width=Pt(0.5))
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text.upper()
    run.font.size = Pt(8)
    run.font.bold = True
    run.font.color.rgb = text_color
    return txBox

def add_card(slide, left, top, width, height,
             fill=CARD, border=BORDA, accent_color=None, accent_left=False):
    r = add_rect(slide, left, top, width, height, fill_color=fill, line_color=border, line_width=Pt(0.5))
    if accent_color and accent_left:
        add_rect(slide, left, top, Pt(3), height, fill_color=accent_color)
    return r

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — CAPA
# ─────────────────────────────────────────────────────────────────────────────
def slide_capa():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, PRETO)

    # Gradiente decorativo (retângulo semi-transparente)
    add_rect(slide, Inches(0), Inches(0), Inches(5), H,
             fill_color=RGBColor(0x12,0x0D,0x02))

    # Linha dourada vertical
    add_rect(slide, Inches(5), Inches(0.5), Pt(3), Inches(6.5), fill_color=OURO)

    # Lion emoji via texto
    add_text(slide, "🦁", Inches(5.3), Inches(0.6), Inches(2), Inches(1.2),
             font_size=72, color=OURO, align=PP_ALIGN.LEFT)

    # Empresa
    add_text(slide, "PASSO A PASSO UNIFORMES  ·  2026",
             Inches(5.3), Inches(1.8), Inches(7.5), Inches(0.4),
             font_size=9, color=CINZA, align=PP_ALIGN.LEFT)

    # Título
    add_text(slide, "Proposta de",
             Inches(5.3), Inches(2.3), Inches(7.5), Inches(0.6),
             font_size=26, bold=False, color=CINZA_CLR, align=PP_ALIGN.LEFT)

    add_text(slide, "Liderança de Vendas",
             Inches(5.3), Inches(2.85), Inches(7.5), Inches(1.0),
             font_size=44, bold=True, color=OURO, align=PP_ALIGN.LEFT)

    add_accent_bar(slide, Inches(5.3), Inches(3.95), Inches(2.5))

    add_text(slide, "Como vamos de R$80k para R$200k em 30 dias\ne construir a cultura que vai durar anos.",
             Inches(5.3), Inches(4.1), Inches(7.5), Inches(0.9),
             font_size=14, color=CINZA_CLR, align=PP_ALIGN.LEFT)

    # Badges
    add_badge(slide, "Abril → Junho 2026", Inches(5.3), Inches(5.3))
    add_badge(slide, "+150% Crescimento",  Inches(7.2), Inches(5.3),
              bg_color=RGBColor(0x05,0x2E,0x1A),
              text_color=VERDE,
              border_color=RGBColor(0x0A,0x5C,0x34))
    add_badge(slide, "3 Novos Canais",     Inches(9.1), Inches(5.3),
              bg_color=RGBColor(0x2B,0x18,0x05),
              text_color=LARANJA,
              border_color=RGBColor(0x78,0x3A,0x06))

    # Info rodapé esquerda
    add_text(slide, "📅  Reunião de Alinhamento · Março 2026\n🏢  Novo Hamburgo, RS    👥  Equipe ~10 pessoas",
             Inches(0.3), Inches(6.4), Inches(4.5), Inches(0.9),
             font_size=9, color=CINZA, align=PP_ALIGN.LEFT)

    # Número do slide
    add_text(slide, "01", Inches(12.8), Inches(7.1), Inches(0.5), Inches(0.35),
             font_size=9, color=CINZA, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — DIAGNÓSTICO
# ─────────────────────────────────────────────────────────────────────────────
def slide_diagnostico():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, PRETO)

    add_badge(slide, "Diagnóstico Honesto", Inches(0.5), Inches(0.35),
              bg_color=RGBColor(0x2B,0x18,0x05), text_color=LARANJA,
              border_color=RGBColor(0x78,0x3A,0x06))
    add_accent_bar(slide, Inches(0.5), Inches(0.82))
    add_text(slide, "Onde estamos hoje", Inches(0.5), Inches(0.95),
             Inches(6), Inches(0.65), font_size=32, bold=True, color=BRANCO)

    # Card faturamento
    add_card(slide, Inches(0.5), Inches(1.75), Inches(2.8), Inches(1.3))
    add_text(slide, "FATURAMENTO ATUAL", Inches(0.6), Inches(1.85),
             Inches(2.6), Inches(0.3), font_size=8, color=CINZA)
    add_text(slide, "R$80k", Inches(0.6), Inches(2.1),
             Inches(2.6), Inches(0.7), font_size=40, bold=True, color=OURO)
    add_text(slide, "por mês · Março 2026", Inches(0.6), Inches(2.75),
             Inches(2.6), Inches(0.25), font_size=9, color=CINZA)

    # Card meta
    add_card(slide, Inches(3.5), Inches(1.75), Inches(2.8), Inches(1.3))
    add_text(slide, "META NOVA", Inches(3.6), Inches(1.85),
             Inches(2.6), Inches(0.3), font_size=8, color=CINZA)
    add_text(slide, "R$200k", Inches(3.6), Inches(2.1),
             Inches(2.6), Inches(0.7), font_size=40, bold=True, color=VERDE)
    add_text(slide, "Abril 2026 · +150% crescimento", Inches(3.6), Inches(2.75),
             Inches(2.6), Inches(0.25), font_size=9, color=CINZA)

    # Problemas
    problemas = [
        ("🔴", "90% dos leads chegam só por WhatsApp + Instagram. Sem diversificação = sem escala."),
        ("🔴", "Sem CRM, leads somem. Sem follow-up estruturado, dinheiro fica na mesa."),
        ("🔴", "Sem cultura de vendas, resultado depende do humor do dia — não de processo."),
        ("🔴", "Mercado Livre, e-commerce e visitas externas = zero explorado."),
    ]
    for i, (emoji, txt) in enumerate(problemas):
        y = Inches(3.25) + i * Inches(0.82)
        add_rect(slide, Inches(0.5), y, Inches(5.8), Inches(0.72),
                 fill_color=RGBColor(0x1A,0x0C,0x0C),
                 line_color=RGBColor(0x4A,0x18,0x18), line_width=Pt(0.5))
        add_rect(slide, Inches(0.5), y, Pt(3), Inches(0.72), fill_color=VERMELHO)
        add_text(slide, emoji + "  " + txt,
                 Inches(0.65), y + Pt(6), Inches(5.55), Inches(0.62),
                 font_size=10.5, color=CINZA_CLR)

    # Box verde lado direito
    add_rect(slide, Inches(6.6), Inches(1.75), Inches(6.4), Inches(5.4),
             fill_color=RGBColor(0x05,0x17,0x10),
             line_color=RGBColor(0x0A,0x3D,0x22), line_width=Pt(0.75))
    add_rect(slide, Inches(6.6), Inches(1.75), Pt(3), Inches(5.4), fill_color=VERDE)

    add_text(slide, "✅  A boa notícia:", Inches(6.8), Inches(2.0),
             Inches(6.0), Inches(0.4), font_size=14, bold=True, color=VERDE)
    add_text(slide,
             "A Passo a Passo tem 30 anos de reputação\n"
             "e qualidade comprovada. A base de clientes\n"
             "é fiel. O produto já é excelente.\n\n"
             "O que falta é uma MÁQUINA DE VENDAS\ntrabalhando por cima disso.",
             Inches(6.8), Inches(2.5), Inches(6.0), Inches(3.5),
             font_size=13, color=CINZA_CLR)

    add_text(slide, "02", Inches(12.8), Inches(7.1), Inches(0.5), Inches(0.35),
             font_size=9, color=CINZA, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — CULTURA LEÕES
# ─────────────────────────────────────────────────────────────────────────────
def slide_cultura():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, PRETO)

    add_rect(slide, Inches(0), Inches(0), Inches(5), H,
             fill_color=RGBColor(0x12,0x0D,0x02))

    add_badge(slide, "A Cultura", Inches(0.5), Inches(0.35),
              bg_color=RGBColor(0x27,0x1E,0x07), text_color=OURO,
              border_color=RGBColor(0x78,0x50,0x1A))
    add_accent_bar(slide, Inches(0.5), Inches(0.82))
    add_text(slide, "🦁 Leões na Selva", Inches(0.5), Inches(0.9),
             Inches(6), Inches(0.75), font_size=34, bold=True, color=OURO)
    add_text(slide, "30 anos de produto. Agora, alma de startup de vendas.",
             Inches(0.5), Inches(1.65), Inches(4.8), Inches(0.4),
             font_size=12, color=CINZA_CLR)

    mandamentos = [
        ("🎯", "Leões caçam — não esperam.", "Prospecção ativa: visitas, ligações, broadcast"),
        ("⚡", "Resposta em menos de 5 min.", "Quem responde primeiro, vende — sempre."),
        ("📋", "CRM antes de fechar o WhatsApp.", "Todo lead cadastrado. Nenhum escapa."),
        ("🤝", "Pede indicação após cada venda.", "Uma venda bem-feita traz outras duas."),
        ("🔔", "Celebra cada venda — o sino toca.", "O time para, vibra junto, placar atualizado."),
    ]
    for i, (emoji, titulo, sub) in enumerate(mandamentos):
        y = Inches(2.2) + i * Inches(0.95)
        add_card(slide, Inches(0.5), y, Inches(4.2), Inches(0.82))
        add_text(slide, emoji, Inches(0.65), y + Pt(8), Inches(0.5), Inches(0.65), font_size=18)
        add_text(slide, titulo, Inches(1.25), y + Pt(6), Inches(3.3), Inches(0.35),
                 font_size=11.5, bold=True, color=BRANCO)
        add_text(slide, sub, Inches(1.25), y + Pt(28), Inches(3.3), Inches(0.4),
                 font_size=9.5, color=CINZA)

    # Pilares — lado direito
    pilares = [
        ("🔔", "O Sino", "Toda venda fechada = sino toca. O valor vai no placar.\nO time vibra junto. Nenhum treinamento substitui isso.", OURO),
        ("📊", "O Placar", "Quadro branco visível: meta do dia (R$7.000),\ntotal realizado, quem vendeu. Competição saudável.", AZUL),
        ("🏆", "Mural dos Campeões", "Na parede da loja: maior venda, hall da fama.\nReconhecimento público vale mais que bônus.", VERDE),
    ]
    for i, (emoji, titulo, descr, cor) in enumerate(pilares):
        y = Inches(1.7) + i * Inches(1.85)
        add_card(slide, Inches(5.3), y, Inches(7.6), Inches(1.7),
                 accent_color=cor, accent_left=True)
        add_text(slide, emoji, Inches(5.6), y + Pt(10), Inches(0.7), Inches(0.5), font_size=26)
        add_text(slide, titulo, Inches(6.4), y + Pt(10), Inches(6.2), Inches(0.4),
                 font_size=13, bold=True, color=cor)
        add_text(slide, descr, Inches(6.4), y + Pt(36), Inches(6.2), Inches(0.8),
                 font_size=10.5, color=CINZA_CLR)

    add_text(slide, "03", Inches(12.8), Inches(7.1), Inches(0.5), Inches(0.35),
             font_size=9, color=CINZA, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — NOVOS CANAIS
# ─────────────────────────────────────────────────────────────────────────────
def slide_canais():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, PRETO)

    add_badge(slide, "Expansão", Inches(0.5), Inches(0.35))
    add_accent_bar(slide, Inches(0.5), Inches(0.82))
    add_text(slide, "3 Novos Canais de Receita", Inches(0.5), Inches(0.9),
             Inches(10), Inches(0.65), font_size=32, bold=True, color=BRANCO)
    add_text(slide, "Hoje: 100% WhatsApp + Instagram.   Em 90 dias: 5 canais gerando receita.",
             Inches(0.5), Inches(1.55), Inches(12), Inches(0.4),
             font_size=12, color=CINZA_CLR)

    canais = [
        {
            "emoji": "🛒", "nome": "Mercado Livre",
            "badge": "GO-LIVE: SEMANA 3 ABR", "badge_cor": LARANJA,
            "badge_bg": RGBColor(0x2B,0x18,0x05),
            "descr": "100 milhões de usuários.\nBusca 'uniforme personalizado'\n= demanda sem atendimento.",
            "accent": OURO,
            "meses": [("Abril", "R$5.000"), ("Maio", "R$12.000"), ("Junho", "R$20.000")],
        },
        {
            "emoji": "🌐", "nome": "E-commerce / Shopify",
            "badge": "GO-LIVE: MAIO", "badge_cor": AZUL,
            "badge_bg": RGBColor(0x0B,0x19,0x30),
            "descr": "Simulador + orçamento.\nSEO local: 'uniforme\nNovo Hamburgo'.",
            "accent": AZUL,
            "meses": [("Abril", "—"), ("Maio", "R$8.000"), ("Junho", "R$20.000")],
        },
        {
            "emoji": "🤝", "nome": "Visitas Externas",
            "badge": "INÍCIO: SEMANA 1 ABR", "badge_cor": VERDE,
            "badge_bg": RGBColor(0x05,0x17,0x10),
            "descr": "Maior ROI de todos.\n5 visitas/semana =\n1 fechamento/semana.",
            "accent": VERDE,
            "meses": [("Abril", "R$14.000"), ("Maio", "R$20.000"), ("Junho", "R$25.000")],
        },
    ]

    for i, c in enumerate(canais):
        x = Inches(0.5) + i * Inches(4.28)
        add_card(slide, x, Inches(2.05), Inches(4.0), Inches(5.1))
        add_rect(slide, x, Inches(2.05), Pt(3), Inches(5.1), fill_color=c["accent"])
        # emoji
        add_text(slide, c["emoji"], x + Inches(0.15), Inches(2.15),
                 Inches(0.7), Inches(0.55), font_size=30)
        # nome
        add_text(slide, c["nome"], x + Inches(0.95), Inches(2.2),
                 Inches(2.9), Inches(0.45), font_size=14, bold=True, color=BRANCO)
        # badge
        add_rect(slide, x + Inches(0.95), Inches(2.65), Inches(2.9), Inches(0.25),
                 fill_color=c["badge_bg"], line_color=c["badge_cor"], line_width=Pt(0.5))
        add_text(slide, c["badge"], x + Inches(0.95), Inches(2.63),
                 Inches(2.9), Inches(0.27), font_size=8, bold=True, color=c["badge_cor"],
                 align=PP_ALIGN.LEFT)
        # descr
        add_text(slide, c["descr"], x + Inches(0.18), Inches(3.05),
                 Inches(3.7), Inches(1.0), font_size=11, color=CINZA_CLR)
        # linha separadora
        add_rect(slide, x + Inches(0.18), Inches(4.15), Inches(3.6), Pt(1), fill_color=BORDA)
        # meses
        for j, (mes, val) in enumerate(c["meses"]):
            yy = Inches(4.25) + j * Inches(0.62)
            add_text(slide, mes, x + Inches(0.18), yy, Inches(1.5), Inches(0.5),
                     font_size=10, color=CINZA)
            cor_val = VERDE if j == 2 else OURO
            add_text(slide, val, x + Inches(2.0), yy, Inches(1.8), Inches(0.5),
                     font_size=10, bold=True, color=cor_val, align=PP_ALIGN.RIGHT)

    # Métricas rodapé
    metricas = [("5 min", "Tempo máximo de resposta"), ("20", "Visitas externas/mês"),
                ("5", "Canais em Junho"), ("30%", "Conversão orçamento→venda")]
    for i, (val, label) in enumerate(metricas):
        x = Inches(0.5) + i * Inches(3.22)
        add_card(slide, x, Inches(7.0), Inches(3.05), Inches(0.45))
        add_text(slide, val, x + Inches(0.15), Inches(7.0), Inches(1.0), Inches(0.45),
                 font_size=16, bold=True, color=OURO)
        add_text(slide, label, x + Inches(1.15), Inches(7.05), Inches(1.8), Inches(0.4),
                 font_size=9, color=CINZA)

    add_text(slide, "04", Inches(12.8), Inches(7.1), Inches(0.5), Inches(0.35),
             font_size=9, color=CINZA, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — METAS E COMISSÕES
# ─────────────────────────────────────────────────────────────────────────────
def slide_comissoes():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, PRETO)

    add_badge(slide, "Motivação", Inches(0.5), Inches(0.35),
              bg_color=RGBColor(0x05,0x17,0x10), text_color=VERDE,
              border_color=RGBColor(0x0A,0x5C,0x34))
    add_accent_bar(slide, Inches(0.5), Inches(0.82), color=VERDE)
    add_text(slide, "Metas e Comissões", Inches(0.5), Inches(0.9),
             Inches(10), Inches(0.65), font_size=32, bold=True, color=BRANCO)
    add_text(slide, "Vendedor que sabe exatamente o que ganha por cada venda, vende mais.",
             Inches(0.5), Inches(1.55), Inches(12), Inches(0.4),
             font_size=12, color=CINZA_CLR)

    # Tabela comissões
    add_text(slide, "COMISSÃO POR CANAL", Inches(0.5), Inches(2.1),
             Inches(6), Inches(0.35), font_size=9, bold=True, color=OURO)
    tabela = [
        ("Canal", "Comissão", True),
        ("Venda direta (WhatsApp / Visita)", "3,5%", False),
        ("Lead Meta Ads convertido", "3,0%", False),
        ("Indicação de cliente antigo ⭐", "4,0%", False),
        ("Mercado Livre", "2,0%", False),
        ("Renovação / Recorrente ⭐", "4,0%", False),
    ]
    for j, (col1, col2, header) in enumerate(tabela):
        yy = Inches(2.45) + j * Inches(0.58)
        bg_c = RGBColor(0x1A,0x13,0x03) if header else CARD
        add_rect(slide, Inches(0.5), yy, Inches(5.8), Inches(0.55),
                 fill_color=bg_c, line_color=BORDA, line_width=Pt(0.4))
        tc1 = OURO if header else BRANCO
        tc2 = OURO if header else VERDE
        add_text(slide, col1, Inches(0.65), yy + Pt(5), Inches(4.5), Inches(0.48),
                 font_size=10.5 if not header else 9,
                 bold=header, color=tc1)
        add_text(slide, col2, Inches(5.0), yy + Pt(5), Inches(1.2), Inches(0.48),
                 font_size=11, bold=True, color=tc2, align=PP_ALIGN.CENTER)

    # Aviso
    add_rect(slide, Inches(0.5), Inches(6.05), Inches(5.8), Inches(0.55),
             fill_color=RGBColor(0x1A,0x0C,0x0C),
             line_color=RGBColor(0x4A,0x18,0x18), line_width=Pt(0.5))
    add_text(slide, "⚠️  Comissão paga após entrega — não no fechamento.",
             Inches(0.65), Inches(6.1), Inches(5.6), Inches(0.45),
             font_size=9.5, color=CINZA_CLR)

    # Bônus lado direito
    add_text(slide, "BÔNUS MENSAIS", Inches(6.8), Inches(2.1),
             Inches(6), Inches(0.35), font_size=9, bold=True, color=OURO)

    bonus = [
        ("🥇", "Maior venda do mês", "R$500 + foto no Mural dos Campeões"),
        ("📦", "Mais pedidos novos", "R$300 + 1 dia de folga"),
        ("🎯", "Equipe bate meta coletiva", "Churrasco pago pela empresa!"),
        ("🤝", "Indicação que fecha", "R$100 extra por indicação convertida"),
        ("💎", "Zero cancelamentos", "R$200 — qualidade recompensada"),
    ]
    for j, (emoji, titulo, descr) in enumerate(bonus):
        yy = Inches(2.55) + j * Inches(0.92)
        add_card(slide, Inches(6.8), yy, Inches(6.1), Inches(0.8))
        add_text(slide, emoji, Inches(6.95), yy + Pt(8), Inches(0.5), Inches(0.65), font_size=20)
        add_text(slide, titulo, Inches(7.6), yy + Pt(6), Inches(3.5), Inches(0.35),
                 font_size=11, bold=True, color=BRANCO)
        add_text(slide, descr, Inches(7.6), yy + Pt(26), Inches(5.1), Inches(0.35),
                 font_size=9.5, color=VERDE)

    add_text(slide, "05", Inches(12.8), Inches(7.1), Inches(0.5), Inches(0.35),
             font_size=9, color=CINZA, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — PROJEÇÃO FINANCEIRA
# ─────────────────────────────────────────────────────────────────────────────
def slide_financeiro():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, PRETO)

    add_badge(slide, "Resultados Projetados", Inches(0.5), Inches(0.35),
              bg_color=RGBColor(0x05,0x17,0x10), text_color=VERDE,
              border_color=RGBColor(0x0A,0x5C,0x34))
    add_accent_bar(slide, Inches(0.5), Inches(0.82), color=VERDE)
    add_text(slide, "Crescimento em 3 Meses", Inches(0.5), Inches(0.9),
             Inches(10), Inches(0.65), font_size=32, bold=True, color=BRANCO)

    # Colunas dos meses
    meses = [
        ("MARÇO", "atual", CINZA,  "R$80k",  False),
        ("ABRIL",  "mês 1", OURO,  "R$200k", True),
        ("MAIO",   "mês 2", CINZA, "R$220k", False),
        ("JUNHO",  "mês 3", VERDE, "R$250k", True),
    ]
    col_x = [Inches(1.3), Inches(4.0), Inches(7.2), Inches(10.0)]
    col_w = Inches(2.5)

    for i, (nome, sub, cor, total, destaque) in enumerate(meses):
        x = col_x[i]
        bc = RGBColor(0x27,0x1E,0x07) if destaque and cor==OURO else \
             (RGBColor(0x05,0x17,0x10) if destaque and cor==VERDE else CARD)
        lc = cor if destaque else BORDA
        add_rect(slide, x, Inches(1.75), col_w, Inches(1.0),
                 fill_color=bc, line_color=lc, line_width=Pt(0.75))
        add_text(slide, nome, x + Inches(0.1), Inches(1.8), col_w - Inches(0.2),
                 Inches(0.35), font_size=9, bold=True, color=cor, align=PP_ALIGN.CENTER)
        add_text(slide, sub,  x + Inches(0.1), Inches(2.15), col_w - Inches(0.2),
                 Inches(0.25), font_size=8, color=CINZA, align=PP_ALIGN.CENTER)

    # Linhas de dados
    linhas = [
        ("WhatsApp direto",    ["R$80k", "R$130k", "R$130k", "R$130k"]),
        ("Meta Ads → CRM",     ["R$0",   "R$40k",  "R$45k",  "R$45k"]),
        ("Visitas externas",   ["R$0",   "R$14k",  "R$20k",  "R$25k"]),
        ("ML + E-commerce",    ["R$0",   "R$5k",   "R$20k",  "R$40k"]),
        ("Indicações",         ["R$0",   "R$11k",  "R$5k",   "R$10k"]),
    ]
    totais_linha = ["R$80k", "R$200k", "R$220k", "R$250k"]

    for j, (label, vals) in enumerate(linhas):
        yy = Inches(2.95) + j * Inches(0.55)
        add_rect(slide, Inches(0), yy, Inches(1.2), Inches(0.52),
                 fill_color=CARD, line_color=BORDA, line_width=Pt(0.3))
        add_text(slide, label, Inches(0.1), yy + Pt(4), Inches(1.1), Inches(0.4),
                 font_size=8.5, color=CINZA_CLR)
        for i, (val) in enumerate(vals):
            x = col_x[i]
            add_rect(slide, x, yy, col_w, Inches(0.52),
                     fill_color=CARD, line_color=BORDA, line_width=Pt(0.3))
            add_text(slide, val, x, yy + Pt(4), col_w, Inches(0.4),
                     font_size=10, color=CINZA_CLR, align=PP_ALIGN.CENTER)

    # Linha TOTAL
    yy = Inches(5.7)
    add_rect(slide, Inches(0), yy, Inches(1.2), Inches(0.6),
             fill_color=ESCURO, line_color=BORDA, line_width=Pt(0.4))
    add_text(slide, "TOTAL", Inches(0.1), yy + Pt(6), Inches(1.1), Inches(0.5),
             font_size=10, bold=True, color=BRANCO)
    for i, (nome, sub, cor, total, destaque) in enumerate(meses):
        x = col_x[i]
        bc = RGBColor(0x27,0x1E,0x07) if cor==OURO and destaque else \
             (RGBColor(0x05,0x17,0x10) if cor==VERDE and destaque else ESCURO)
        add_rect(slide, x, yy, col_w, Inches(0.6),
                 fill_color=bc, line_color=cor if destaque else BORDA, line_width=Pt(0.75))
        add_text(slide, total, x, yy + Pt(6), col_w, Inches(0.5),
                 font_size=14, bold=True, color=cor, align=PP_ALIGN.CENTER)

    # 3 métricas finais
    metricas_fin = [
        ("R$670k", "total 3 meses", OURO),
        ("R$12,2k", "investimento total", VERDE),
        ("55x", "retorno s/ investimento", LARANJA),
    ]
    for i, (val, label, cor) in enumerate(metricas_fin):
        x = Inches(0.5) + i * Inches(4.28)
        add_card(slide, x, Inches(6.55), Inches(4.0), Inches(0.75))
        add_text(slide, val, x + Inches(0.15), Inches(6.6),
                 Inches(1.8), Inches(0.65), font_size=28, bold=True, color=cor)
        add_text(slide, label, x + Inches(2.1), Inches(6.72),
                 Inches(1.8), Inches(0.45), font_size=9, color=CINZA)

    add_text(slide, "06", Inches(12.8), Inches(7.1), Inches(0.5), Inches(0.35),
             font_size=9, color=CINZA, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — RITMO DIÁRIO
# ─────────────────────────────────────────────────────────────────────────────
def slide_ritmo():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, PRETO)

    add_badge(slide, "Processo", Inches(0.5), Inches(0.35),
              bg_color=RGBColor(0x0B,0x19,0x30), text_color=AZUL,
              border_color=RGBColor(0x1E,0x40,0xAF))
    add_accent_bar(slide, Inches(0.5), Inches(0.82), color=AZUL)
    add_text(slide, "A Rotina do Leão", Inches(0.5), Inches(0.9),
             Inches(8), Inches(0.65), font_size=32, bold=True, color=BRANCO)
    add_text(slide, "Resultado consistente vem de rotina consistente. Todo dia, sem exceção.",
             Inches(0.5), Inches(1.55), Inches(12), Inches(0.4),
             font_size=12, color=CINZA_CLR)

    timeline = [
        ("08:00", "Reunião de Abertura (15 min)",
         "Resultado ontem · meta hoje · campanha do dia · técnica da semana", OURO),
        ("08:15", "Ataque Matinal",
         "Todos os leads do Kommo respondidos. WhatsApps pendentes resolvidos.", LARANJA),
        ("11:00", "Check-in Rápido (5 min)",
         "Orçamentos enviados? Alguém precisa de ajuda para fechar?", AZUL),
        ("14:00", "Campanha da Tarde",
         "Broadcast WhatsApp. Stories com oferta/urgência no Instagram.", VERDE),
        ("17:30", "Fechamento do Dia",
         "Atualizar Kommo. Registrar no placar. Planejar amanhã.", OURO),
    ]

    for i, (hora, titulo, descr, cor) in enumerate(timeline):
        yy = Inches(2.1) + i * Inches(1.0)
        # Dot
        add_rect(slide, Inches(0.5), yy + Inches(0.08), Inches(0.15), Inches(0.15),
                 fill_color=cor)
        # Conector
        if i < len(timeline)-1:
            add_rect(slide, Inches(0.565), yy + Inches(0.23), Pt(1.5), Inches(0.75),
                     fill_color=BORDA)
        # Hora
        add_text(slide, hora, Inches(0.8), yy, Inches(0.85), Inches(0.4),
                 font_size=10, bold=True, color=cor)
        # Card
        add_card(slide, Inches(1.75), yy, Inches(4.5), Inches(0.82))
        add_text(slide, titulo, Inches(1.9), yy + Pt(5), Inches(4.2), Inches(0.35),
                 font_size=11, bold=True, color=BRANCO)
        add_text(slide, descr, Inches(1.9), yy + Pt(26), Inches(4.2), Inches(0.4),
                 font_size=9.5, color=CINZA)

    # Campanhas por dia — lado direito
    add_text(slide, "CAMPANHAS POR DIA DA SEMANA", Inches(6.8), Inches(2.1),
             Inches(6), Inches(0.35), font_size=9, bold=True, color=OURO)

    campanhas = [
        ("Segunda", "URGÊNCIA", "Prazo fechando — garanta sua produção essa semana", LARANJA),
        ("Terça",   "PROVA SOCIAL", "Foto de cliente satisfeito + depoimento real", VERDE),
        ("Quarta",  "OFERTA", "Pedidos +20 peças: frete grátis no Vale dos Sinos", OURO),
        ("Quinta",  "EDUCAÇÃO", "Sabia que dá pra personalizar até o número na manga?", AZUL),
        ("Sexta",   "FECHAMENTO", "Última chance de garantir produção essa semana!", LARANJA),
    ]
    for j, (dia, badge_txt, texto, cor) in enumerate(campanhas):
        yy = Inches(2.55) + j * Inches(0.9)
        add_card(slide, Inches(6.8), yy, Inches(6.1), Inches(0.78))
        add_text(slide, dia, Inches(6.95), yy + Pt(5), Inches(0.9), Inches(0.65),
                 font_size=11, bold=True, color=BRANCO)
        add_rect(slide, Inches(7.95), yy + Pt(8), Inches(1.3), Inches(0.22),
                 fill_color=ESCURO, line_color=cor, line_width=Pt(0.5))
        add_text(slide, badge_txt, Inches(7.95), yy + Pt(6), Inches(1.3), Inches(0.24),
                 font_size=7, bold=True, color=cor, align=PP_ALIGN.CENTER)
        add_text(slide, texto, Inches(9.4), yy + Pt(5), Inches(3.4), Inches(0.65),
                 font_size=9.5, color=CINZA_CLR)

    add_text(slide, "07", Inches(12.8), Inches(7.1), Inches(0.5), Inches(0.35),
             font_size=9, color=CINZA, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — O QUE EU PRECISO
# ─────────────────────────────────────────────────────────────────────────────
def slide_preciso():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, PRETO)

    add_badge(slide, "Para Executar", Inches(0.5), Inches(0.35))
    add_accent_bar(slide, Inches(0.5), Inches(0.82))
    add_text(slide, "O que preciso da liderança", Inches(0.5), Inches(0.9),
             Inches(10), Inches(0.65), font_size=32, bold=True, color=BRANCO)
    add_text(slide, "Itens simples, baixo custo. Retorno: 55x sobre o investimento.",
             Inches(0.5), Inches(1.55), Inches(12), Inches(0.4),
             font_size=12, color=CINZA_CLR)

    add_text(slide, "INFRAESTRUTURA", Inches(0.5), Inches(2.1),
             Inches(6), Inches(0.35), font_size=9, bold=True, color=OURO)

    infra = [
        ("🔔", "Sino físico para a loja",        "Cowbell ou campainha", "R$50",   VERDE),
        ("📊", "Quadro branco + marcadores",       "Placar diário de vendas", "R$150", VERDE),
        ("💻", "Kommo CRM ativo",                  "Já em configuração · ~R$400/mês", "Em andamento", OURO),
        ("📸", "Sessão de fotos dos produtos",     "10 produtos para ML e Shopify", "R$400",  LARANJA),
        ("📣", "Meta Ads: R$3.000 em Abril",       "Dobrar para escalar leads no lançamento", "R$3.000", LARANJA),
    ]
    for j, (emoji, titulo, sub, valor, cor) in enumerate(infra):
        yy = Inches(2.5) + j * Inches(0.82)
        add_card(slide, Inches(0.5), yy, Inches(5.8), Inches(0.72))
        add_text(slide, emoji, Inches(0.65), yy + Pt(8), Inches(0.5), Inches(0.6), font_size=18)
        add_text(slide, titulo, Inches(1.25), yy + Pt(6), Inches(3.5), Inches(0.32),
                 font_size=11, bold=True, color=BRANCO)
        add_text(slide, sub, Inches(1.25), yy + Pt(26), Inches(3.5), Inches(0.32),
                 font_size=9, color=CINZA)
        add_text(slide, valor, Inches(5.0), yy + Pt(10), Inches(1.2), Inches(0.45),
                 font_size=10, bold=True, color=cor, align=PP_ALIGN.CENTER)

    # Autoridade
    add_text(slide, "AUTORIDADE NECESSÁRIA", Inches(6.8), Inches(2.1),
             Inches(6), Inches(0.35), font_size=9, bold=True, color=VERDE)

    autoridade = [
        ("✅", "Autonomia para visitas externas",
         "Autorização para visitar escolas e empresas em nome da Passo a Passo"),
        ("✅", "Aprovação das comissões",
         "Definir os percentuais finais junto com a diretoria antes de Abril"),
        ("✅", "Reunião de kickoff com equipe",
         "01/Abril — apresentar o plano, instalar o sino, definir metas"),
        ("✅", "Acesso ao Mercado Livre e Shopify",
         "Criar e gerenciar os canais digitais em nome da empresa"),
    ]
    for j, (emoji, titulo, descr) in enumerate(autoridade):
        yy = Inches(2.55) + j * Inches(1.05)
        add_card(slide, Inches(6.8), yy, Inches(6.1), Inches(0.88))
        add_text(slide, emoji + "  " + titulo,
                 Inches(6.95), yy + Pt(5), Inches(5.8), Inches(0.35),
                 font_size=11, bold=True, color=BRANCO)
        add_text(slide, descr, Inches(6.95), yy + Pt(28), Inches(5.8), Inches(0.45),
                 font_size=9.5, color=CINZA)

    # Box resumo
    add_rect(slide, Inches(6.8), Inches(6.7), Inches(6.1), Inches(0.65),
             fill_color=RGBColor(0x27,0x1E,0x07),
             line_color=RGBColor(0x78,0x50,0x1A), line_width=Pt(0.75))
    add_text(slide, "Investimento total Abril: ~R$4.200   →   Retorno esperado: R$200.000",
             Inches(6.95), Inches(6.75), Inches(5.8), Inches(0.5),
             font_size=10.5, bold=True, color=OURO, align=PP_ALIGN.CENTER)

    add_text(slide, "08", Inches(12.8), Inches(7.1), Inches(0.5), Inches(0.35),
             font_size=9, color=CINZA, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — FECHAMENTO
# ─────────────────────────────────────────────────────────────────────────────
def slide_fechamento():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, PRETO)

    # Fundo decorativo
    add_rect(slide, Inches(0), Inches(0), W, H,
             fill_color=RGBColor(0x09,0x08,0x02))
    add_rect(slide, Inches(4), Inches(0), Inches(5.33), H,
             fill_color=RGBColor(0x0F,0x0C,0x02))

    add_text(slide, "🦁", Inches(5.5), Inches(0.3), Inches(2.5), Inches(2.0),
             font_size=96, color=OURO, align=PP_ALIGN.CENTER)

    add_text(slide, "A Passo a Passo tem",
             Inches(2.0), Inches(2.3), Inches(9.0), Inches(0.7),
             font_size=28, bold=False, color=CINZA_CLR, align=PP_ALIGN.CENTER)

    add_text(slide, "30 anos de produto.",
             Inches(2.0), Inches(3.0), Inches(9.0), Inches(0.75),
             font_size=36, bold=True, color=OURO, align=PP_ALIGN.CENTER)

    add_text(slide, "Agora vamos ter uma máquina de vendas.",
             Inches(2.0), Inches(3.75), Inches(9.0), Inches(0.65),
             font_size=28, bold=True, color=LARANJA, align=PP_ALIGN.CENTER)

    add_accent_bar(slide, Inches(5.8), Inches(4.55), Inches(1.8), color=OURO)

    add_text(slide, "Não precisamos de sorte.\nPrecisamos de processo, disciplina e cultura.\nEu estou pronto para liderar essa transformação.",
             Inches(2.5), Inches(4.7), Inches(8.0), Inches(1.0),
             font_size=13, color=CINZA_CLR, align=PP_ALIGN.CENTER)

    # Badges
    badges_f = [("R$200k em Abril", OURO, RGBColor(0x27,0x1E,0x07), RGBColor(0x78,0x50,0x1A)),
                ("R$670k em 3 meses", VERDE, RGBColor(0x05,0x17,0x10), RGBColor(0x0A,0x5C,0x34)),
                ("5 Canais em Junho", AZUL, RGBColor(0x0B,0x19,0x30), RGBColor(0x1E,0x40,0xAF))]
    total_w = sum([Inches(2.0)] * 3) + Inches(0.2) * 2
    start_x = (W - total_w) / 2
    for i, (label, tc, bg_c, bc) in enumerate(badges_f):
        x = start_x + i * Inches(2.2)
        add_rect(slide, x, Inches(5.85), Inches(2.0), Inches(0.3),
                 fill_color=bg_c, line_color=bc, line_width=Pt(0.5))
        add_text(slide, label.upper(), x, Inches(5.85), Inches(2.0), Inches(0.3),
                 font_size=8, bold=True, color=tc, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(3.5), Inches(6.4), Inches(6.3), Inches(0.75),
             fill_color=RGBColor(0x27,0x1E,0x07),
             line_color=RGBColor(0x78,0x50,0x1A), line_width=Pt(1))
    add_text(slide, "Quando começamos?  →  01 de Abril de 2026",
             Inches(3.5), Inches(6.45), Inches(6.3), Inches(0.65),
             font_size=15, bold=True, color=OURO, align=PP_ALIGN.CENTER)

    add_text(slide, "09", Inches(12.8), Inches(7.1), Inches(0.5), Inches(0.35),
             font_size=9, color=CINZA, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# Gerar todos os slides
# ─────────────────────────────────────────────────────────────────────────────
slide_capa()
slide_diagnostico()
slide_cultura()
slide_canais()
slide_comissoes()
slide_financeiro()
slide_ritmo()
slide_preciso()
slide_fechamento()

# Salvar
output = r"C:\Users\gusta\OneDrive\Área de Trabalho\PassoaPasso\docs\apresentacao-lider-vendas-2026.pptx"
prs.save(output)
print(f"PowerPoint salvo em: {output}")
