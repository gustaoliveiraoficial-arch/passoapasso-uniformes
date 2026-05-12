# -*- coding: utf-8 -*-
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Cores Passo a Passo ──────────────────────────────────────────────────────
NAVY      = RGBColor(0x0A, 0x16, 0x28)
NAVY2     = RGBColor(0x0D, 0x1F, 0x3C)
NAVY3     = RGBColor(0x1A, 0x2E, 0x4A)
LARANJA   = RGBColor(0xEA, 0x58, 0x0C)
LARANJA2  = RGBColor(0xF9, 0x73, 0x16)
LARANJA3  = RGBColor(0xFB, 0x92, 0x3C)
BRANCO    = RGBColor(0xFF, 0xFF, 0xFF)
CINZA     = RGBColor(0x94, 0xA3, 0xB8)
CINZA_CLR = RGBColor(0xE2, 0xE8, 0xF0)
VERDE     = RGBColor(0x10, 0xB9, 0x81)
VERMELHO  = RGBColor(0xEF, 0x44, 0x44)
AMARELO   = RGBColor(0xF5, 0x9E, 0x0B)
AZUL      = RGBColor(0x38, 0xBD, 0xF8)
CARD      = RGBColor(0x0F, 0x28, 0x45)
BORDA     = RGBColor(0x1E, 0x3A, 0x5F)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── helpers ──────────────────────────────────────────────────────────────────
def bg(slide, color=NAVY):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def R(slide, l, t, w, h, fill=None, line=None, lw=Pt(0.5)):
    s = slide.shapes.add_shape(1, l, t, w, h)
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if line: s.line.color.rgb = line; s.line.width = lw
    else: s.line.fill.background()
    return s

def T(slide, text, l, t, w, h, size=12, bold=False, color=BRANCO,
      align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return tb

def bar(slide, l, t, w=Inches(0.4), color=LARANJA):
    R(slide, l, t, w, Pt(3), fill=color)

def card(slide, l, t, w, h, fill=CARD, line=BORDA, accent=None):
    R(slide, l, t, w, h, fill=fill, line=line, lw=Pt(0.5))
    if accent: R(slide, l, t, Pt(3), h, fill=accent)

def badge(slide, txt, l, t, bg_c=LARANJA, tc=BRANCO, w=Inches(2.0)):
    R(slide, l, t, w, Inches(0.28), fill=bg_c)
    T(slide, txt.upper(), l, t, w, Inches(0.28), size=8, bold=True,
      color=tc, align=PP_ALIGN.CENTER)

def sim(slide, l, t): # check verde
    R(slide, l, t, Inches(0.24), Inches(0.24),
      fill=RGBColor(0x05,0x28,0x16), line=RGBColor(0x10,0xB9,0x81), lw=Pt(0.5))
    T(slide, "OK", l, t+Pt(1), Inches(0.24), Inches(0.22),
      size=6.5, bold=True, color=VERDE, align=PP_ALIGN.CENTER)

def nao(slide, l, t):
    R(slide, l, t, Inches(0.24), Inches(0.24),
      fill=RGBColor(0x28,0x05,0x05), line=RGBColor(0xEF,0x44,0x44), lw=Pt(0.5))
    T(slide, "NAO", l, t+Pt(1), Inches(0.24), Inches(0.22),
      size=5.5, bold=True, color=VERMELHO, align=PP_ALIGN.CENTER)

def aviso(slide, l, t):
    R(slide, l, t, Inches(0.24), Inches(0.24),
      fill=RGBColor(0x28,0x1C,0x05), line=AMARELO, lw=Pt(0.5))
    T(slide, "!", l, t+Pt(1), Inches(0.24), Inches(0.22),
      size=8, bold=True, color=AMARELO, align=PP_ALIGN.CENTER)

def num(slide, n):
    T(slide, str(n).zfill(2), Inches(12.85), Inches(7.1),
      Inches(0.45), Inches(0.35), size=9, color=CINZA, align=PP_ALIGN.RIGHT)

def header(slide, titulo, subtitulo="", badge_txt="", badge_cor=LARANJA):
    R(slide, Inches(0), Inches(0), W, Inches(1.05), fill=NAVY2)
    R(slide, Inches(0), Inches(0), W, Pt(4), fill=LARANJA)
    if badge_txt:
        badge(slide, badge_txt, Inches(0.4), Inches(0.22), bg_c=badge_cor, w=Inches(2.0))
    T(slide, titulo, Inches(0.4) if not badge_txt else Inches(2.55),
      Inches(0.1), Inches(10), Inches(0.55),
      size=28, bold=True, color=BRANCO)
    if subtitulo:
        T(slide, subtitulo, Inches(0.4), Inches(0.65), Inches(12.5), Inches(0.38),
          size=11, color=CINZA_CLR)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ═════════════════════════════════════════════════════════════════════════════
def s1():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, NAVY)
    R(slide, Inches(0), Inches(0), Inches(0.08), H, fill=LARANJA)
    R(slide, Inches(0), H-Pt(4), W, Pt(4), fill=LARANJA)
    R(slide, Inches(0), Inches(0), W, Pt(4), fill=LARANJA)
    R(slide, Inches(0.08), Inches(0), Inches(4.5), H, fill=NAVY2)
    T(slide, "ACADEMIA\nDE PRODUTO", Inches(0.3), Inches(0.9), Inches(4.1),
      Inches(2.0), size=36, bold=True, color=LARANJA, align=PP_ALIGN.CENTER)
    T(slide, "Passo a Passo\nUniformes", Inches(0.3), Inches(3.0), Inches(4.1),
      Inches(1.0), size=18, color=CINZA_CLR, align=PP_ALIGN.CENTER)
    R(slide, Inches(0.4), Inches(4.2), Inches(3.7), Pt(2), fill=LARANJA)
    T(slide, "2026", Inches(0.3), Inches(4.35), Inches(4.1), Inches(0.45),
      size=12, color=CINZA, align=PP_ALIGN.CENTER)
    T(slide, "Tudo que voce precisa saber\npara vender com confianca.",
      Inches(4.9), Inches(1.2), Inches(8.0), Inches(1.1),
      size=28, bold=True, color=BRANCO)
    T(slide, "Nesta aula voce vai aprender:",
      Inches(4.9), Inches(2.5), Inches(8.0), Inches(0.4),
      size=13, color=CINZA_CLR)
    topicos = ["Todos os nossos tecidos — o que sao, como funcionam, para quem indicar",
               "Serigrafia, DTF, Bordado e Sublimacao — quando e como usar cada um",
               "Guia pratico: qual estampa vai em qual tecido",
               "Argumentos de venda por segmento de cliente"]
    for i, t in enumerate(topicos):
        y = Inches(2.95) + i * Inches(0.72)
        R(slide, Inches(4.9), y+Pt(8), Inches(0.06), Inches(0.3), fill=LARANJA)
        T(slide, t, Inches(5.1), y, Inches(7.7), Inches(0.65),
          size=12, color=CINZA_CLR)
    num(slide, 1)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — MAPA DE PRODUTOS
# ═════════════════════════════════════════════════════════════════════════════
def s2():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, NAVY)
    header(slide, "O que a Passo a Passo vende?",
           "Conheca todas as nossas linhas antes de aprofundar nos tecidos.", "VISAO GERAL")
    linhas = [
        (LARANJA,  "CAMISETAS",    "Poliviscose, Algodao 381,\nAlgodao PA, Dry-Fit JIMP,\nDry-Fit Liso, Dry-Fit Manchester"),
        (AZUL,     "POLO",         "Polo Poliviscose\nPolo Dry-Fit JIMP\nPolo Dry Manchester\nPolo Piquet PV"),
        (VERDE,    "SOCIAL",       "Camisa Social\n(lancamento)\nManga bufante\nJaleco"),
        (AMARELO,  "FITNESS",      "Legging\nShort / Bermuda\nTop / Cropped\nCrop Tradicional"),
        (LARANJA2, "INVERNO",      "Jaqueta M.1201\nCorta Vento\nThermoblock\nMoletom / Japona\nColete"),
        (CINZA,    "OUTROS",       "Avental\nBody Infantil\nSacolas\nManga longa"),
    ]
    for i, (cor, nome, items) in enumerate(linhas):
        col = i % 3; row = i // 2
        x = Inches(0.35) + col * Inches(4.33)
        y = Inches(1.3) + row * Inches(2.85)
        card(slide, x, y, Inches(4.1), Inches(2.6), accent=cor)
        R(slide, x, y, Inches(4.1), Pt(3), fill=cor)
        T(slide, nome, x+Inches(0.2), y+Pt(8), Inches(3.7), Inches(0.4),
          size=14, bold=True, color=cor)
        T(slide, items, x+Inches(0.2), y+Pt(34), Inches(3.7), Inches(2.0),
          size=10.5, color=CINZA_CLR)
    num(slide, 2)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — TECIDO POLIVISCOSE
# ═════════════════════════════════════════════════════════════════════════════
def s3():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, NAVY)
    header(slide, "Poliviscose (PV)", "67% poliester + 33% viscose", "TECIDO #1")
    # Bloco info
    infos = [
        ("O que e?", "Mistura de poliester (durabilidade) com viscose (maciez e frescor). Muito usado em uniformes corporativos."),
        ("Caracteristicas", "Toque fresco. Seca rapido. Nao enruga facilmente. Duravel. Superficie lisa e uniforme — otima para estampas."),
        ("Para quem indicar", "Escolas, empresas, comercio — qualquer segmento que quer qualidade com preco acessivel."),
        ("Argumento de venda", "\"Dura mais que algodao, nao amassa, seca rapido. Perfeito para uso profissional diario.\""),
    ]
    for i, (titulo, descr) in enumerate(infos):
        y = Inches(1.2) + i * Inches(1.38)
        card(slide, Inches(0.35), y, Inches(7.2), Inches(1.22), accent=LARANJA)
        T(slide, titulo, Inches(0.6), y+Pt(6), Inches(6.8), Inches(0.32),
          size=10, bold=True, color=LARANJA)
        T(slide, descr, Inches(0.6), y+Pt(26), Inches(6.8), Inches(0.85),
          size=11, color=CINZA_CLR)
    # Estampas lado direito
    T(slide, "ESTAMPAS", Inches(7.85), Inches(1.2), Inches(5.1), Inches(0.3),
      size=9, bold=True, color=CINZA)
    estampas = [
        ("Serigrafia SILK", VERDE, "Excelente — tinta adere perfeitamente"),
        ("DTF", VERDE, "Excelente — cores vibrantes em qualquer cor"),
        ("Bordado", VERDE, "Bom — tecido firme facilita"),
        ("Sublimacao", VERMELHO, "NAO funciona — tem viscose (fibra natural)"),
    ]
    for j, (nome, cor, desc) in enumerate(estampas):
        y = Inches(1.6) + j * Inches(1.35)
        card(slide, Inches(7.85), y, Inches(5.1), Inches(1.18),
             fill=CARD if cor==VERDE else RGBColor(0x1A,0x06,0x06))
        R(slide, Inches(7.85), y, Pt(3), Inches(1.18), fill=cor)
        T(slide, nome, Inches(8.1), y+Pt(6), Inches(4.6), Inches(0.35),
          size=12, bold=True, color=cor)
        T(slide, desc, Inches(8.1), y+Pt(28), Inches(4.6), Inches(0.7),
          size=10, color=CINZA_CLR)
    num(slide, 3)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ALGODOES (381 e PA)
# ═════════════════════════════════════════════════════════════════════════════
def s4():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, NAVY)
    header(slide, "Algodao 381 e Algodao PA",
           "100% algodao — o classico que nunca sai de moda", "TECIDO #2 e #3")
    # Algodao 381
    T(slide, "ALGODAO 381", Inches(0.35), Inches(1.25),
      Inches(6.15), Inches(0.35), size=9, bold=True, color=LARANJA)
    a381 = [
        ("Fio fino e denso, toque macio e premium. Mais pesado que malha comum."),
        ("Ideal para uniformes de alta qualidade — corporativos, formandos, academia."),
        ("Cores: Preto, Branco, Off-White, Azul Marinho (10 pcs). Demais cores: 90 pcs."),
        ("Argumento: \"Aquela camiseta que parece cara, que todo mundo quer tocar.\""),
    ]
    for i, t in enumerate(a381):
        y = Inches(1.65) + i * Inches(0.75)
        R(slide, Inches(0.38), y+Pt(7), Inches(0.06), Inches(0.25), fill=LARANJA)
        T(slide, t, Inches(0.55), y, Inches(5.8), Inches(0.65), size=10.5, color=CINZA_CLR)
    # Algodao PA
    T(slide, "ALGODAO PA (PENTEADO)", Inches(6.7), Inches(1.25),
      Inches(6.2), Inches(0.35), size=9, bold=True, color=AZUL)
    apa = [
        ("Algodao penteado = fibras curtas removidas. Resultado: fio mais resistente e macio."),
        ("30% mais resistente que algodao comum. Dura muito mais."),
        ("Disponivel em mais cores: Preto, Mescla claro/escuro, Verde, Rosa, Laranja neon."),
        ("Argumento: \"Camiseta que nao bola, nao afina, nao perde a cor tao rapido.\""),
    ]
    for i, t in enumerate(apa):
        y = Inches(1.65) + i * Inches(0.75)
        R(slide, Inches(6.73), y+Pt(7), Inches(0.06), Inches(0.25), fill=AZUL)
        T(slide, t, Inches(6.9), y, Inches(6.0), Inches(0.65), size=10.5, color=CINZA_CLR)
    # Linha divisoria
    R(slide, Inches(6.55), Inches(1.2), Pt(1), Inches(4.0), fill=BORDA)
    # Estampas — linha de baixo
    R(slide, Inches(0.35), Inches(5.5), Inches(12.6), Pt(1), fill=BORDA)
    T(slide, "ESTAMPAS PARA ALGODAO:", Inches(0.35), Inches(5.6), Inches(3.5), Inches(0.4),
      size=10, bold=True, color=CINZA)
    ests = [("Serigrafia", VERDE, "Perfeita"), ("DTF", VERDE, "Muito bom"),
            ("Bordado", VERDE, "Perfeito"), ("Sublimacao", VERMELHO, "NAO — so poliester")]
    for i, (nome, cor, desc) in enumerate(ests):
        x = Inches(3.5) + i * Inches(2.4)
        card(slide, x, Inches(5.55), Inches(2.2), Inches(0.85))
        R(slide, x, Inches(5.55), Inches(2.2), Pt(3), fill=cor)
        T(slide, nome, x+Inches(0.1), Inches(5.62), Inches(2.0), Inches(0.3),
          size=10, bold=True, color=cor)
        T(slide, desc, x+Inches(0.1), Inches(5.88), Inches(2.0), Inches(0.4),
          size=9, color=CINZA_CLR)
    num(slide, 4)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — DRY-FIT (JIMP, LISO, MANCHESTER)
# ═════════════════════════════════════════════════════════════════════════════
def s5():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, NAVY)
    header(slide, "Linha Dry-Fit",
           "100% poliester — tecnologia que mantém a pessoa seca mesmo suando", "TECIDOS #4, 5 e 6")
    # Explicacao base
    card(slide, Inches(0.35), Inches(1.2), Inches(12.6), Inches(1.05),
         fill=RGBColor(0x08,0x1E,0x3A))
    R(slide, Inches(0.35), Inches(1.2), Inches(12.6), Pt(3), fill=AZUL)
    T(slide, "Como funciona o Dry-Fit?", Inches(0.55), Inches(1.25),
      Inches(12.0), Inches(0.32), size=11, bold=True, color=AZUL)
    T(slide, "Fibras de poliester sao ocas por dentro. Quando a pessoa sua, o suor e transportado de dentro para fora do tecido e evapora rapido. Resultado: a pessoa se sente seca e confortavel mesmo se exercitando ou trabalhando em ambiente quente.",
      Inches(0.55), Inches(1.55), Inches(12.0), Inches(0.6), size=10.5, color=CINZA_CLR)
    # 3 tipos
    tipos = [
        (LARANJA, "DRY-FIT JIMP", "Textura JIMP (pontos em relevo)",
         "O mais tecnico. Textura aumenta a evaporacao do suor. Mais moderno e esportivo.",
         "Esportes, academias, empresas jovens/dinamicas, equipes de campo."),
        (VERDE,  "DRY-FIT LISO",  "Superficie lisa (sem textura)",
         "Mais elegante que o JIMP. Aparencia limpa e profissional. Leve e confortavel.",
         "Uniformes corporativos, escolas, empresas de tecnologia, equipes de vendas."),
        (AZUL,   "DRY-FIT MANCHESTER", "Trama semi-translucida / polo",
         "Mais sofisticado. Parece um polo mas com tecnologia dry-fit. Respirabilidade extra.",
         "Uniformes executivos premium, polo de alta qualidade, segmentos formais."),
    ]
    for i, (cor, nome, sub, desc, ind) in enumerate(tipos):
        x = Inches(0.35) + i * Inches(4.33)
        card(slide, x, Inches(2.45), Inches(4.1), Inches(4.8), accent=cor)
        T(slide, nome, x+Inches(0.22), Inches(2.55), Inches(3.7), Inches(0.4),
          size=12, bold=True, color=cor)
        T(slide, sub, x+Inches(0.22), Inches(2.92), Inches(3.7), Inches(0.32),
          size=9, italic=True, color=CINZA)
        R(slide, x+Inches(0.22), Inches(3.28), Inches(3.5), Pt(1), fill=cor)
        T(slide, desc, x+Inches(0.22), Inches(3.38), Inches(3.7), Inches(1.35),
          size=10, color=CINZA_CLR)
        T(slide, "Indicado:", x+Inches(0.22), Inches(4.72), Inches(3.7), Inches(0.28),
          size=8, bold=True, color=cor)
        T(slide, ind, x+Inches(0.22), Inches(4.98), Inches(3.7), Inches(0.85),
          size=9.5, color=CINZA_CLR)
        # Estampas rapidas
        T(slide, "Sublimacao: OK  |  DTF: OK  |  Serigrafia: OK  |  Bordado: com cuidado",
          x+Inches(0.22), Inches(5.98), Inches(3.7), Inches(0.35),
          size=7.5, color=CINZA)
    num(slide, 5)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PIQUET PV + VISCOLYCRA
# ═════════════════════════════════════════════════════════════════════════════
def s6():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, NAVY)
    header(slide, "Piquet PV e Viscolycra",
           "Tecidos com diferenciais especificos — conheca quando indicar", "TECIDOS #7 e #8")
    # Piquet PV
    T(slide, "PIQUET PV", Inches(0.35), Inches(1.28),
      Inches(6.2), Inches(0.32), size=9, bold=True, color=LARANJA)
    card(slide, Inches(0.35), Inches(1.62), Inches(6.2), Inches(5.45), accent=LARANJA)
    piquet = [
        ("Composicao", "67% poliester + 33% viscose com TRATAMENTO ANTIPILLING."),
        ("O que e Antipilling?", "Antipilling = nao forma bolinhas. As camisas usadas por meses continuam novas."),
        ("Textura Pique", "Trama diagonal 3D visivel — aparencia premium, elegante. Diferencial visual imediato."),
        ("Toque", "Fresco como o poliviscose, mas com textura exclusiva que valoriza o uniforme."),
        ("Quando indicar", "Polos corporativos de alta qualidade. Clientes que querem diferencial visual."),
        ("Argumento de venda", "\"Este tecido nao forma bolinhas. Depois de 1 ano de uso, a camisa continua bonita.\""),
    ]
    for i, (tit, desc) in enumerate(piquet):
        y = Inches(1.72) + i * Inches(0.82)
        T(slide, tit + ":", Inches(0.55), y, Inches(1.9), Inches(0.3),
          size=9, bold=True, color=LARANJA)
        T(slide, desc, Inches(2.55), y, Inches(3.8), Inches(0.7),
          size=10, color=CINZA_CLR)
    # Viscolycra
    T(slide, "VISCOLYCRA", Inches(6.8), Inches(1.28),
      Inches(6.1), Inches(0.32), size=9, bold=True, color=AZUL)
    card(slide, Inches(6.8), Inches(1.62), Inches(6.1), Inches(5.45), accent=AZUL)
    viscol = [
        ("Composicao", "96% viscose + 4% elastano (lycra)."),
        ("O que faz?", "Viscose da maciez e frescor. Elastano da elasticidade. Tecido ajusta ao corpo."),
        ("Toque", "Premium, fluido, levemente brilhante. Parece caro ao toque."),
        ("Elasticidade", "Estica e volta. Nao deforma com o uso. Mantém o caimento."),
        ("Quando indicar", "Manga bufante feminina, uniformes ajustados, setores criativos."),
        ("Cuidado estampas", "Bordado dificil (estica). Serigrafia so com tinta elastica especial. DTF: bom resultado."),
    ]
    for i, (tit, desc) in enumerate(viscol):
        y = Inches(1.72) + i * Inches(0.82)
        T(slide, tit + ":", Inches(7.0), y, Inches(1.9), Inches(0.3),
          size=9, bold=True, color=AZUL)
        T(slide, desc, Inches(9.0), y, Inches(3.7), Inches(0.7),
          size=10, color=CINZA_CLR)
    num(slide, 6)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — TECIDOS DE INVERNO
# ═════════════════════════════════════════════════════════════════════════════
def s7():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, NAVY)
    header(slide, "Tecidos da Linha Inverno",
           "Thermoblock, Corta Vento, Moletom, Matelasse e Brim Sarja", "TECIDOS #9 a #13")
    inv = [
        (LARANJA2, "THERMOBLOCK", "Poliester + camada termica interna",
         "Jaquetas e coletes. Mantem o calor corporal sem pesar. Ideal para RS no inverno e ambientes com ar-condicionado.",
         "Serigrafia: bom | DTF: bom | Bordado: excelente"),
        (AZUL, "CORTA VENTO", "Nylon (poliamida) impermeavel",
         "Repele vento e chuva. Levissimo, pode ser dobrado em um bolso. Nao abafa como japona pesada.",
         "DTF: possivel | Serigrafia: dificil | Bordado: nao recomendado"),
        (VERDE, "MOLETOM", "Algodao ou poliester (fleece/french terry)",
         "College, canguru, basico. Macio por dentro, estruturado por fora. Quente e confortavel.",
         "Serigrafia: excelente | DTF: bom | Bordado: excelente"),
        (AMARELO, "MATELASSE", "Multiplas camadas acolchoadas",
         "Coletes premium. Textura 3D exclusiva. Aparencia sofisticada e executiva. Alta durabilidade.",
         "Bordado: dificil | DTF: possivel | Serigrafia: dificil"),
        (LARANJA, "BRIM SARJA", "100% algodao industrial diagonal",
         "Bermudas, calcas industriais. Resistencia extrema. Nao derrete. Para trabalho pesado.",
         "Serigrafia: perfeita | DTF: bom | Bordado: perfeito"),
    ]
    for i, (cor, nome, comp, desc, est) in enumerate(inv):
        col = i % 3 if i < 3 else (i - 3)
        row = 0 if i < 3 else 1
        x = Inches(0.35) + (col if i < 3 else (col+0.65)) * Inches(4.33)
        if i >= 3: x = Inches(2.25) + (i-3) * Inches(4.7)
        y = Inches(1.28) + row * Inches(2.95)
        w = Inches(4.1) if i < 3 else Inches(4.5)
        card(slide, x, y, w, Inches(2.7), accent=cor)
        T(slide, nome, x+Inches(0.2), y+Pt(6), w-Inches(0.3), Inches(0.35),
          size=12, bold=True, color=cor)
        T(slide, comp, x+Inches(0.2), y+Pt(28), w-Inches(0.3), Inches(0.28),
          size=8, italic=True, color=CINZA)
        R(slide, x+Inches(0.2), y+Pt(52), w-Inches(0.4), Pt(1), fill=cor)
        T(slide, desc, x+Inches(0.2), y+Pt(60), w-Inches(0.3), Inches(1.2),
          size=9.5, color=CINZA_CLR)
        T(slide, est, x+Inches(0.2), y+Inches(2.3), w-Inches(0.3), Inches(0.32),
          size=7.5, color=CINZA)
    num(slide, 7)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — SERIGRAFIA SILK
# ═════════════════════════════════════════════════════════════════════════════
def s8():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, NAVY)
    header(slide, "Serigrafia SILK (Silk Screen)",
           "A tecnica mais classica — tinta que penetra nas fibras do tecido", "ESTAMPA #1")
    # Como funciona
    card(slide, Inches(0.35), Inches(1.18), Inches(6.2), Inches(3.1))
    T(slide, "COMO FUNCIONA (simples):", Inches(0.55), Inches(1.28),
      Inches(5.8), Inches(0.3), size=9, bold=True, color=LARANJA)
    passos = ["1. Cria-se uma TELA (malha com emulsao fotossensivel)",
              "2. O logo/arte e gravado na tela com luz UV",
              "3. A tela e colocada sobre o tecido",
              "4. Tinta e passada por cima com um RODO",
              "5. Tinta atravessa os furos e penetra nas fibras",
              "6. Tinta seca e fica permanente — resiste a muitas lavagens"]
    for i, p in enumerate(passos):
        T(slide, p, Inches(0.55), Inches(1.62)+i*Inches(0.4),
          Inches(5.9), Inches(0.38), size=10, color=CINZA_CLR)
    # Vantagens
    card(slide, Inches(0.35), Inches(4.42), Inches(6.2), Inches(2.72), accent=VERDE)
    T(slide, "VANTAGENS:", Inches(0.55), Inches(4.52),
      Inches(5.8), Inches(0.3), size=9, bold=True, color=VERDE)
    vantagens = ["Cores VIBRANTES — tinta dentro da fibra, nao na superficie",
                 "Durabilidade excepcional — aguenta 50+ lavagens sem desbotar",
                 "Custo baixissimo por peca em grande volume (50+ pcs iguais)",
                 "Funciona em quase todos os tecidos"]
    for i, v in enumerate(vantagens):
        T(slide, "• " + v, Inches(0.55), Inches(4.86)+i*Inches(0.45),
          Inches(5.9), Inches(0.4), size=10.5, color=CINZA_CLR)
    # Limitacoes
    card(slide, Inches(6.8), Inches(1.18), Inches(6.1), Inches(2.35),
         fill=RGBColor(0x1A,0x0A,0x0A), accent=VERMELHO)
    T(slide, "LIMITACOES (saiba explicar):", Inches(7.0), Inches(1.28),
      Inches(5.7), Inches(0.3), size=9, bold=True, color=VERMELHO)
    lims = ["Cada COR diferente = uma tela separada (custo sobe)",
            "Designs com muitas cores ficam caros em poucas pecas",
            "Gradientes suaves sao dificeis de reproduzir",
            "Setup inicial demora (gravacao das telas)"]
    for i, l in enumerate(lims):
        T(slide, "• " + l, Inches(7.0), Inches(1.62)+i*Inches(0.48),
          Inches(5.8), Inches(0.42), size=10, color=CINZA_CLR)
    # Quando indicar
    card(slide, Inches(6.8), Inches(3.65), Inches(6.1), Inches(1.55), accent=LARANJA)
    T(slide, "QUANDO INDICAR:", Inches(7.0), Inches(3.75),
      Inches(5.7), Inches(0.3), size=9, bold=True, color=LARANJA)
    T(slide, "Logo simples (1-3 cores) + pedido grande (50+ pcs iguais)\nCliente quer DURABILIDADE maxima e preco baixo por peca",
      Inches(7.0), Inches(4.08), Inches(5.8), Inches(0.9), size=11, color=CINZA_CLR)
    # Tecidos compatíveis
    card(slide, Inches(6.8), Inches(5.3), Inches(6.1), Inches(1.85))
    T(slide, "FUNCIONA EM:", Inches(7.0), Inches(5.38),
      Inches(5.7), Inches(0.3), size=9, bold=True, color=VERDE)
    T(slide, "Algodao 381  Algodao PA  Poliviscose\nPiquet PV  Dry-Fit (todos)  Moletom  Brim Sarja\nCorta vento (dificil)  Matelasse (dificil)",
      Inches(7.0), Inches(5.7), Inches(5.8), Inches(1.2), size=10.5, color=CINZA_CLR)
    num(slide, 8)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — DTF
# ═════════════════════════════════════════════════════════════════════════════
def s9():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, NAVY)
    header(slide, "DTF — Direct to Film",
           "A revolucao das estampas: funciona em quase tudo, com qualquer design", "ESTAMPA #2")
    card(slide, Inches(0.35), Inches(1.18), Inches(6.2), Inches(3.35))
    T(slide, "COMO FUNCIONA (simples):", Inches(0.55), Inches(1.28),
      Inches(5.8), Inches(0.3), size=9, bold=True, color=AZUL)
    passos = ["1. Design e impresso em um FILME PET especial (papel filme)",
              "2. Po adesivo de poliamida e espalhado sobre o filme impresso",
              "3. Filme vai ao forno para ativar o adesivo (po derrete ligeiramente)",
              "4. Filme e posicionado sobre o tecido",
              "5. PRENSA TERMICA a 160C por alguns segundos",
              "6. Filme e removido — estampa fica permanentemente no tecido"]
    for i, p in enumerate(passos):
        T(slide, p, Inches(0.55), Inches(1.62)+i*Inches(0.45),
          Inches(5.9), Inches(0.4), size=10, color=CINZA_CLR)
    card(slide, Inches(0.35), Inches(4.65), Inches(6.2), Inches(2.5), accent=AZUL)
    T(slide, "POR QUE O DTF E INCRIVEL:", Inches(0.55), Inches(4.75),
      Inches(5.8), Inches(0.3), size=9, bold=True, color=AZUL)
    ptq = ["Funciona em QUALQUER cor de tecido (claro ou escuro)",
           "Gradientes, sombras, fotos — qualquer design complexo",
           "Sem preparacao do tecido (economiza tempo e dinheiro)",
           "Toque macio — estampa nao fica borrachuda",
           "Funciona em algodao, dry-fit, poliviscose e mistos"]
    for i, p in enumerate(ptq):
        T(slide, "• " + p, Inches(0.55), Inches(5.1)+i*Inches(0.4),
          Inches(5.9), Inches(0.35), size=10.5, color=CINZA_CLR)
    # Comparação com Silk
    card(slide, Inches(6.8), Inches(1.18), Inches(6.1), Inches(3.1))
    T(slide, "DTF vs SERIGRAFIA — QUANDO USAR CADA UM:", Inches(7.0), Inches(1.28),
      Inches(5.7), Inches(0.4), size=9, bold=True, color=CINZA_CLR)
    linhas = [("SITUACAO", "DTF", "SERIGRAFIA", True),
              ("Design colorido (4+ cores)", "✓ Perfeito", "✗ Caro", False),
              ("Design simples (1-2 cores)", "OK", "✓ Melhor preco", False),
              ("Pedido pequeno (10-30 pcs)", "✓ Melhor", "✗ Setup caro", False),
              ("Pedido grande (100+ pcs iguais)", "OK", "✓ Mais barato/peca", False),
              ("Tecido escuro", "✓ Funciona", "✓ Funciona", False),
              ("Gradiente/Foto", "✓ Perfeito", "✗ Dificil", False)]
    for i, (sit, d, s, hdr) in enumerate(linhas):
        y = Inches(1.72) + i*Inches(0.5)
        fc = RGBColor(0x14,0x25,0x3A) if hdr else (CARD if i%2==0 else NAVY2)
        R(slide, Inches(6.8), y, Inches(6.1), Inches(0.47), fill=fc, line=BORDA, lw=Pt(0.3))
        c1 = CINZA if hdr else CINZA_CLR
        c2 = VERDE if ("✓" in d) else (VERMELHO if "✗" in d else CINZA_CLR)
        c3 = VERDE if ("✓" in s) else (VERMELHO if "✗" in s else CINZA_CLR)
        T(slide, sit, Inches(6.95), y+Pt(5), Inches(2.7), Inches(0.38), size=9.5, bold=hdr, color=c1)
        T(slide, d, Inches(9.7), y+Pt(5), Inches(1.4), Inches(0.38), size=9.5, bold=hdr, color=c2, align=PP_ALIGN.CENTER)
        T(slide, s, Inches(11.15), y+Pt(5), Inches(1.6), Inches(0.38), size=9.5, bold=hdr, color=c3, align=PP_ALIGN.CENTER)
    # Tecidos
    card(slide, Inches(6.8), Inches(5.3), Inches(6.1), Inches(1.85))
    T(slide, "FUNCIONA EM QUASE TUDO:", Inches(7.0), Inches(5.38),
      Inches(5.7), Inches(0.3), size=9, bold=True, color=VERDE)
    T(slide, "Algodao  Dry-Fit  Poliviscose  Piquet PV  Viscolycra\nMoletom  Thermoblock  Brim Sarja — praticamente tudo!\nCorta vento e matelasse: possivel com cuidado",
      Inches(7.0), Inches(5.7), Inches(5.8), Inches(1.2), size=10.5, color=CINZA_CLR)
    num(slide, 9)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — BORDADO
# ═════════════════════════════════════════════════════════════════════════════
def s10():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, NAVY)
    header(slide, "Bordado em Maquina",
           "A estampa mais premium e duravel — aparencia executiva garantida", "ESTAMPA #3")
    card(slide, Inches(0.35), Inches(1.18), Inches(6.2), Inches(3.2))
    T(slide, "COMO FUNCIONA:", Inches(0.55), Inches(1.28),
      Inches(5.8), Inches(0.3), size=9, bold=True, color=LARANJA)
    passos = ["1. Logo e convertido para arquivo de bordado digital (programado em ponto a ponto)",
              "2. Tecido e fixado em um BASTIDOR para nao escorregar",
              "3. Maquina automatica move agulha e bastidor simultaneamente",
              "4. Linhas coloridas sao aplicadas ponto a ponto sobre o tecido",
              "5. Resultado: desenho costurado — nao e tinta, e fio de linha"]
    for i, p in enumerate(passos):
        T(slide, p, Inches(0.55), Inches(1.62)+i*Inches(0.5),
          Inches(5.9), Inches(0.44), size=10, color=CINZA_CLR)
    card(slide, Inches(0.35), Inches(4.5), Inches(6.2), Inches(2.65), accent=LARANJA)
    T(slide, "POR QUE INDICAR BORDADO:", Inches(0.55), Inches(4.6),
      Inches(5.8), Inches(0.3), size=9, bold=True, color=LARANJA)
    vant = ["Aparencia EXECUTIVA e premium — parece logo de empresa grande",
            "Durabilidade extrema — linha costurada nao sai nunca",
            "Tato diferente — logo em relevo, da para sentir",
            "Consistencia perfeita — todas as pecas identicas",
            "Ideal para LOGOS CORPORATIVOS e uniformes de alto padrao"]
    for i, v in enumerate(vant):
        T(slide, "• " + v, Inches(0.55), Inches(4.98)+i*Inches(0.42),
          Inches(5.9), Inches(0.38), size=10.5, color=CINZA_CLR)
    # Limitacoes e indicacoes
    card(slide, Inches(6.8), Inches(1.18), Inches(6.1), Inches(2.35),
         fill=RGBColor(0x1A,0x0A,0x0A), accent=VERMELHO)
    T(slide, "LIMITACOES:", Inches(7.0), Inches(1.28),
      Inches(5.7), Inches(0.3), size=9, bold=True, color=VERMELHO)
    lims = ["Pecas menores — bordado ocupa espaco limitado (logo grande fica caro)",
            "Tecidos finos (dry-fit muito fino) — agulha pode danificar",
            "Tecidos muito elasticos (viscolycra) — deforma durante o processo",
            "Nao faz gradientes ou fotos — apenas cores solidas e linhas"]
    for i, l in enumerate(lims):
        T(slide, "• " + l, Inches(7.0), Inches(1.62)+i*Inches(0.48),
          Inches(5.8), Inches(0.42), size=10, color=CINZA_CLR)
    card(slide, Inches(6.8), Inches(3.65), Inches(6.1), Inches(1.5), accent=VERDE)
    T(slide, "QUANDO INDICAR BORDADO:", Inches(7.0), Inches(3.75),
      Inches(5.7), Inches(0.3), size=9, bold=True, color=VERDE)
    T(slide, "Logos corporativos  •  Uniformes executivos\nEscolhas de tecidos firmes: Piquet, Algodao, Moletom, Brim\nCliente que quer aparencia PREMIUM",
      Inches(7.0), Inches(4.08), Inches(5.8), Inches(0.9), size=10.5, color=CINZA_CLR)
    card(slide, Inches(6.8), Inches(5.3), Inches(6.1), Inches(1.85))
    T(slide, "MELHORES TECIDOS PARA BORDADO:", Inches(7.0), Inches(5.38),
      Inches(5.7), Inches(0.3), size=9, bold=True, color=LARANJA)
    T(slide, "Piquet PV: PERFEITO  |  Algodao: PERFEITO  |  Moletom: OTIMO\nBrim Sarja: PERFEITO  |  Poliviscose: BOM  |  Thermoblock: BOM\nDry-Fit fino: COM CUIDADO  |  Viscolycra: DIFICIL  |  Corta Vento: NAO",
      Inches(7.0), Inches(5.7), Inches(5.8), Inches(1.25), size=10, color=CINZA_CLR)
    num(slide, 10)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — SUBLIMACAO
# ═════════════════════════════════════════════════════════════════════════════
def s11():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, NAVY)
    header(slide, "Sublimacao",
           "A tecnica dos uniformes esportivos — cores que nao saem nunca", "ESTAMPA #4")
    card(slide, Inches(0.35), Inches(1.18), Inches(6.2), Inches(3.1))
    T(slide, "COMO FUNCIONA:", Inches(0.55), Inches(1.28),
      Inches(5.8), Inches(0.3), size=9, bold=True, color=VERDE)
    passos = ["1. Arte e impressa em papel especial com TINTA SUBLIMATICA",
              "2. Papel e posicionado sobre o tecido poliester",
              "3. PRENSA TERMICA a 180-200C por 10-15 segundos",
              "4. O CALOR faz a tinta mudar de solido direto para GAS (sublimacao)",
              "5. Gas penetra nas fibras do tecido e se solidifica DENTRO delas",
              "6. Cor fica INTEGRADA a fibra — impossivel sair em lavagens"]
    for i, p in enumerate(passos):
        T(slide, p, Inches(0.55), Inches(1.62)+i*Inches(0.42),
          Inches(5.9), Inches(0.38), size=10, color=CINZA_CLR)
    # Explicacao algodao
    card(slide, Inches(0.35), Inches(4.4), Inches(6.2), Inches(2.75),
         fill=RGBColor(0x18,0x08,0x02), accent=VERMELHO)
    T(slide, "POR QUE NAO FUNCIONA EM ALGODAO? (IMPORTANTE SABER!)", Inches(0.55), Inches(4.5),
      Inches(5.8), Inches(0.35), size=9, bold=True, color=VERMELHO)
    T(slide, "Poliester e uma fibra TERMOPL ASTICA: quando aquece a 180C, as fibras\n\"abrem\" e a tinta gasosa consegue entrar.\n\nAlgodao e uma fibra NATURAL: nao se abre com calor. A tinta gasosa\nnao consegue penetrar. Resultado: estampa apagada e sai na 1a lavagem.\n\nCLIENTE PERGUNTAR: \"Da pra fazer sublimacao em algodao?\"\nRESPOSTA: \"Nao. Para estampa colorida em algodao, recomendo DTF.\"",
      Inches(0.55), Inches(4.88), Inches(5.9), Inches(2.0), size=10, color=CINZA_CLR)
    # Vantagens
    card(slide, Inches(6.8), Inches(1.18), Inches(6.1), Inches(2.55), accent=VERDE)
    T(slide, "VANTAGENS DA SUBLIMACAO:", Inches(7.0), Inches(1.28),
      Inches(5.7), Inches(0.3), size=9, bold=True, color=VERDE)
    vant = ["Cores EXTREMAMENTE vibrantes e saturadas",
            "Estampa full body — cobre TODA a peca",
            "Toque normal — nao fica emborrachado",
            "Durabilidade extrema — nao desbota nunca",
            "Ideal para uniformes esportivos com muito detalhe"]
    for i, v in enumerate(vant):
        T(slide, "• " + v, Inches(7.0), Inches(1.65)+i*Inches(0.42),
          Inches(5.8), Inches(0.38), size=10.5, color=CINZA_CLR)
    # Regra de ouro
    card(slide, Inches(6.8), Inches(3.85), Inches(6.1), Inches(1.5),
         fill=RGBColor(0x12,0x1E,0x08), accent=AMARELO)
    T(slide, "REGRA DE OURO:", Inches(7.0), Inches(3.95),
      Inches(5.7), Inches(0.3), size=10, bold=True, color=AMARELO)
    T(slide, "Sublimacao APENAS em tecidos 100% sinteticos.\nPoliester = perfeito. Dry-Fit = perfeito.\nAlgodao = NUNCA. Mistos = resultado fraco.",
      Inches(7.0), Inches(4.3), Inches(5.8), Inches(0.85), size=11, color=CINZA_CLR)
    card(slide, Inches(6.8), Inches(5.45), Inches(6.1), Inches(1.7))
    T(slide, "TECIDOS PARA SUBLIMACAO:", Inches(7.0), Inches(5.55),
      Inches(5.7), Inches(0.3), size=9, bold=True, color=VERDE)
    T(slide, "Dry-Fit JIMP: PERFEITO  |  Dry-Fit Liso: PERFEITO\nDry-Fit Manchester: MUITO BOM  |  Corta Vento 100% poliester: BOM\nALGODAO (qualquer tipo): NAO FUNCIONA",
      Inches(7.0), Inches(5.88), Inches(5.8), Inches(1.1), size=10.5, color=CINZA_CLR)
    num(slide, 11)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — MATRIZ TECIDO x ESTAMPA
# ═════════════════════════════════════════════════════════════════════════════
def s12():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, NAVY)
    header(slide, "Guia Rapido: Qual estampa vai em qual tecido?",
           "Cole este guia na sua mesa — consulte antes de cada atendimento", "TABELA DE COMPATIBILIDADE")
    # Tabela
    tecidos = ["Poliviscose", "Algodao 381", "Algodao PA", "Dry-Fit JIMP", "Dry-Fit Liso",
               "Dry Manchester", "Piquet PV", "Viscolycra", "Thermoblock", "Corta Vento",
               "Moletom", "Brim Sarja", "Matelasse"]
    # [Serigrafia, DTF, Bordado, Sublimacao]
    # 0=ok, 1=nao, 2=cuidado
    matriz = [
        [0,0,0,1], [0,0,0,1], [0,0,0,1], [0,0,2,0], [0,0,2,0],
        [0,0,2,0], [0,0,0,1], [2,0,2,1], [0,0,0,2], [2,2,1,2],
        [0,0,0,2], [0,0,0,1], [2,2,1,2],
    ]
    hdrs = ["TECIDO", "SERIGRAFIA", "DTF", "BORDADO", "SUBLIMACAO"]
    col_x = [Inches(0.25), Inches(3.35), Inches(5.85), Inches(8.35), Inches(10.85)]
    col_w = [Inches(3.05), Inches(2.4), Inches(2.4), Inches(2.4), Inches(2.1)]
    # Cabeçalho da tabela
    R(slide, Inches(0.25), Inches(1.2), Inches(12.75), Inches(0.45),
      fill=NAVY3, line=BORDA, lw=Pt(0.4))
    for j, h in enumerate(hdrs):
        c = LARANJA if j==0 else [CINZA,VERDE,VERDE,LARANJA2,VERDE][j]
        T(slide, h, col_x[j], Inches(1.25), col_w[j], Inches(0.38),
          size=8, bold=True, color=c, align=PP_ALIGN.CENTER)
    for i, (tecido, compat) in enumerate(zip(tecidos, matriz)):
        y = Inches(1.68) + i * Inches(0.4)
        fc = CARD if i%2==0 else NAVY2
        R(slide, Inches(0.25), y, Inches(12.75), Inches(0.38), fill=fc, line=BORDA, lw=Pt(0.2))
        T(slide, tecido, col_x[0]+Pt(4), y+Pt(5), col_w[0], Inches(0.3),
          size=9.5, color=CINZA_CLR)
        for j, c in enumerate(compat):
            cx = col_x[j+1] + (col_w[j+1]-Inches(0.24))/2
            if c==0: sim(slide, cx, y+Pt(8))
            elif c==1: nao(slide, cx, y+Pt(8))
            else: aviso(slide, cx, y+Pt(8))
    # Legenda
    ly = Inches(7.05)
    sim(slide, Inches(0.4), ly)
    T(slide, "OK — Funciona bem", Inches(0.72), ly, Inches(2.5), Inches(0.3), size=8.5, color=VERDE)
    aviso(slide, Inches(3.5), ly)
    T(slide, "COM CUIDADO — possivel mas dificil", Inches(3.82), ly, Inches(3.5), Inches(0.3), size=8.5, color=AMARELO)
    nao(slide, Inches(7.5), ly)
    T(slide, "NAO — nao funciona neste tecido", Inches(7.82), ly, Inches(3.5), Inches(0.3), size=8.5, color=VERMELHO)
    num(slide, 12)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — ARGUMENTOS POR SEGMENTO
# ═════════════════════════════════════════════════════════════════════════════
def s13():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, NAVY)
    header(slide, "Argumentos de Venda por Segmento",
           "Use este guia para recomendar o produto certo para cada cliente", "GUIA DE VENDAS")
    segs = [
        (LARANJA, "ESCOLA / ENSINO", "Poliviscose ou Algodao PA",
         "\"Dura o ano letivo inteiro. Nao amassa. Facil de lavar. Cria identidade visual da escola.\"",
         "Silk: logo simples  |  DTF: logo colorido"),
        (AZUL, "EMPRESA / CORPORATIVO", "Piquet PV ou Poliviscose",
         "\"Nao forma bolinhas. Depois de 1 ano de uso, a camisa continua bonita. Passa profissionalismo.\"",
         "Bordado: logo executivo  |  DTF: arte detalhada"),
        (VERDE, "ACADEMIA / ESPORTE", "Dry-Fit JIMP ou Dry-Fit Liso",
         "\"O aluno nao transpira visivel. Seca em minutos. Cores vibrantes com sublimacao que nao saem.\"",
         "Sublimacao: arte total  |  DTF: logo colorido"),
        (AMARELO, "TIME / FUTEBOL", "Dry-Fit Liso ou Dry-Fit Manchester",
         "\"Uniforme personalizado com numero, nome e logo. Leve, confortavel, cores que nao desbotam.\"",
         "Sublimacao total  |  DTF para logos"),
        (LARANJA2, "SAUDE / CLINICA", "Piquet PV ou Algodao PA",
         "\"Aparencia profissional. Facil de higienizar. Nao perde a forma com muitas lavagens.\"",
         "Bordado: logo clinica  |  Silk: texto simples"),
        (CINZA, "INDUSTRIA / CAMPO", "Brim Sarja ou Dry-Fit",
         "\"Resistencia maxima. Brim nao rasga facil, nao derrete. Dry-Fit mantém o trabalhador fresco.\"",
         "Silk: logo simples  |  DTF: arte personalizada"),
    ]
    for i, (cor, nome, tecido, arg, est) in enumerate(segs):
        col = i % 2; row = i // 2
        x = Inches(0.35) + col * Inches(6.5)
        y = Inches(1.28) + row * Inches(2.02)
        card(slide, x, y, Inches(6.2), Inches(1.88), accent=cor)
        T(slide, nome, x+Inches(0.22), y+Pt(6), Inches(5.7), Inches(0.32),
          size=11, bold=True, color=cor)
        T(slide, "Tecido: " + tecido, x+Inches(0.22), y+Pt(28),
          Inches(5.7), Inches(0.28), size=9, italic=True, color=CINZA)
        T(slide, arg, x+Inches(0.22), y+Pt(50), Inches(5.7), Inches(0.75),
          size=10, color=CINZA_CLR)
        T(slide, est, x+Inches(0.22), y+Pt(100), Inches(5.7), Inches(0.25),
          size=8, color=CINZA)
    num(slide, 13)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — REGRAS DE OURO (RESUMAO FINAL)
# ═════════════════════════════════════════════════════════════════════════════
def s14():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, NAVY)
    header(slide, "As 6 Regras de Ouro do Vendedor",
           "Salve, fixe na mesa, repita ate sair de memoria", "RESUMAO FINAL")
    regras = [
        (VERMELHO, "1", "Sublimacao = SO EM DRY-FIT (poliester)",
         "Algodao nao sublima. Se o cliente quer design colorido em algodao — indica DTF."),
        (AZUL, "2", "DTF funciona em QUALQUER TECIDO",
         "Design complexo, cores escuras, gradiente, poucas pecas — DTF e sempre seguro."),
        (LARANJA, "3", "Serigrafia = volume + logo simples",
         "1-3 cores + 50+ pecas iguais = serigrafia ganha no preco. Menos pecas = DTF."),
        (VERDE, "4", "Bordado = logo executivo premium",
         "Quando o cliente quer aparencia de empresa grande. Piquet + bordado = uniforme de alto nivel."),
        (AMARELO, "5", "Piquet PV nao forma bolinhas — use isso como diferencial",
         "\"Antipilling\" e argumento de venda forte. Mostre a diferenca para o cliente."),
        (LARANJA2, "6", "Dry-Fit resolve o problema do calor",
         "Colaborador que trabalha em ambiente quente ou se mexe muito — dry-fit e a indicacao certa."),
    ]
    for i, (cor, n, titulo, desc) in enumerate(regras):
        col = i % 2; row = i // 2
        x = Inches(0.35) + col * Inches(6.5)
        y = Inches(1.28) + row * Inches(2.0)
        card(slide, x, y, Inches(6.2), Inches(1.85))
        R(slide, x, y, Inches(0.6), Inches(1.85), fill=cor)
        T(slide, n, x, y+Pt(18), Inches(0.6), Inches(1.2),
          size=32, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
        T(slide, titulo, x+Inches(0.72), y+Pt(8), Inches(5.3), Inches(0.38),
          size=12, bold=True, color=cor)
        T(slide, desc, x+Inches(0.72), y+Pt(34), Inches(5.3), Inches(1.1),
          size=10.5, color=CINZA_CLR)
    num(slide, 14)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — FECHAMENTO
# ═════════════════════════════════════════════════════════════════════════════
def s15():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, NAVY)
    R(slide, Inches(0), Inches(0), W, Pt(4), fill=LARANJA)
    R(slide, Inches(0), H-Pt(4), W, Pt(4), fill=LARANJA)
    R(slide, Inches(0), Inches(0), Inches(0.08), H, fill=LARANJA)
    T(slide, "Agora voce conhece",
      Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.75),
      size=40, bold=False, color=CINZA_CLR, align=PP_ALIGN.CENTER)
    T(slide, "o que voce vende.",
      Inches(0.5), Inches(2.25), Inches(12.3), Inches(0.85),
      size=50, bold=True, color=LARANJA, align=PP_ALIGN.CENTER)
    R(slide, Inches(5.7), Inches(3.2), Inches(2.0), Pt(3), fill=LARANJA)
    T(slide, "Conhecimento sobre produto gera CONFIANCA.\nConfianca fecha VENDAS.\nVoce nao esta mais so mostrando camiseta — voce esta consultando.",
      Inches(1.5), Inches(3.4), Inches(10.3), Inches(1.2),
      size=15, italic=True, color=CINZA_CLR, align=PP_ALIGN.CENTER)
    itens = ["Conhece cada tecido e para quem indicar",
             "Sabe qual estampa vai em qual tecido",
             "Tem argumentos de venda por segmento",
             "Pode responder qualquer pergunta tecnica do cliente"]
    for i, it in enumerate(itens):
        x = Inches(0.5) + (i % 2) * Inches(6.4)
        y = Inches(4.8) + (i // 2) * Inches(0.75)
        R(slide, x, y+Pt(7), Inches(0.08), Inches(0.25), fill=LARANJA)
        T(slide, it, x+Inches(0.2), y, Inches(6.0), Inches(0.65), size=13, color=CINZA_CLR)
    T(slide, "Passo a Passo Uniformes  |  Academia de Produto  |  2026",
      Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.3),
      size=9, color=CINZA, align=PP_ALIGN.CENTER)
    num(slide, 15)

# ── Gerar ────────────────────────────────────────────────────────────────────
s1(); s2(); s3(); s4(); s5(); s6(); s7()
s8(); s9(); s10(); s11(); s12(); s13(); s14(); s15()

out = r"C:\Users\gusta\OneDrive\Área de Trabalho\PassoaPasso\docs\aula-produtos-tecidos-2026.pptx"
prs.save(out)
print(f"Salvo: {out}")
